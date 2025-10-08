from flask import Flask, redirect, url_for, render_template, request, flash, session, send_from_directory, jsonify, abort
from functools import wraps
from flask_sqlalchemy import SQLAlchemy
from flask_bcrypt import Bcrypt
from flask_migrate import Migrate
from flask_mail import Mail, Message as MailMessage
from itsdangerous import URLSafeTimedSerializer, SignatureExpired, BadSignature
from datetime import datetime, timedelta, timezone
import os
import json
import csv
import re
import requests
import random
import io
from werkzeug.utils import secure_filename
from sqlalchemy import or_, func, case
from config import get_config

# For PPT parsing
from pptx import Presentation

# For PDF parsing
import fitz  # PyMuPDF


app = Flask(__name__)

# Load configuration from config.py based on environment
app.config.from_object(get_config())


def utcnow() -> datetime:
    """Return the current UTC datetime with timezone awareness."""
    return datetime.now(timezone.utc)


def _ensure_utc(dt: datetime | None) -> datetime | None:
    """Return a timezone-aware UTC datetime for stored timestamps."""
    if dt is None:
        return None
    if dt.tzinfo is None:
        return dt.replace(tzinfo=timezone.utc)
    return dt.astimezone(timezone.utc)

# -------------------- Database Setup -------------------- #
app.config['STUDENT_HUB_UPLOAD_FOLDER'] = os.path.join(app.root_path, 'static', 'uploads', 'student_hub')

ALLOWED_STUDENT_HUB_EXTENSIONS = {"pdf", "pptx", "docx", "png", "jpg", "jpeg"}

os.makedirs(app.config['STUDENT_HUB_UPLOAD_FOLDER'], exist_ok=True)
app.config['SESSION_PERMANENT'] = False
app.config['SESSION_TYPE'] = 'filesystem'
db = SQLAlchemy(app)
bcrypt = Bcrypt(app)
migrate = Migrate(app, db)  # ✅ NEW
mail = Mail(app)  # ✅ Flask-Mail for password reset emails

# Initialize password reset token serializer
serializer = URLSafeTimedSerializer(app.config['SECRET_KEY'])

# -------------------- User Model -------------------- #
# -------------------- User Model -------------------- #
class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(150), unique=True, nullable=False)
    password = db.Column(db.String(200), nullable=False)
    full_name = db.Column(db.String(200))
    age = db.Column(db.Integer)
    phone_number = db.Column(db.String(50))
    email = db.Column(db.String(120), unique=True, nullable=False)
    role = db.Column(db.String(20), default="user")
    created_at = db.Column(db.DateTime, default=utcnow)
    last_login = db.Column(db.DateTime)

    # Password reset fields
    reset_token = db.Column(db.String(200), nullable=True)
    reset_token_expiry = db.Column(db.DateTime, nullable=True)

    # Terms of Service acceptance
    accepted_terms = db.Column(db.Boolean, default=False, nullable=False)
    terms_accepted_at = db.Column(db.DateTime, nullable=True)

    # ✅ Courses allowed (comma-separated list for now)
    courses = db.Column(db.String(500), default="seerah,arabic")

    course_progress = db.relationship(
        'UserCourseProgress',
        backref='user',
        lazy=True,
        cascade="all, delete-orphan"
    )

    student_files = db.relationship(
        'StudentHubFile',
        backref='user',
        lazy=True,
        cascade="all, delete-orphan"
    )

    quiz_attempts = db.relationship(
        'QuizAttempt',
        backref='user',
        lazy=True,
        cascade="all, delete-orphan"
    )

    quiz_retake_attempts = db.relationship(
        'QuizRetakeAttempt',
        backref='user',
        lazy=True,
        cascade="all, delete-orphan"
    )

    exam_attempts = db.relationship(
        'ExamAttempt',
        backref='user',
        lazy=True,
        cascade="all, delete-orphan"
    )

    course_agreements = db.relationship(
        'CourseAgreement',
        backref='user',
        lazy=True,
        cascade="all, delete-orphan"
    )

    course_access = db.relationship(
        'CourseAccess',
        backref='user',
        lazy=True,
        cascade="all, delete-orphan"
    )

    def get_courses(self):
        return self.courses.split(",") if self.courses else []

    def set_courses(self, course_list):
        self.courses = ",".join(course_list)

    def has_course_access(self, course_id):
        """Check if user has access to a specific course via new access control system"""
        course = Course.query.get(course_id)
        if not course:
            return False
        return course.user_has_access(self.id)
# -------------------- Course & Lesson Models -------------------- #
class Course(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    year = db.Column(db.Integer, nullable=False)
    description = db.Column(db.Text)

    # ✅ NEW: Hierarchy & monetization fields
    course_type = db.Column(db.String(20), default='standalone')  # standalone, year_based, sub_course
    parent_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=True)
    order_index = db.Column(db.Integer, default=0)
    is_published = db.Column(db.Boolean, default=True)  # ✅ Changed to True by default for backward compatibility
    price = db.Column(db.Numeric(10, 2), nullable=True)  # nullable = free
    show_on_homepage = db.Column(db.Boolean, default=False)  # ✅ Show course on homepage
    course_assignment = db.Column(db.String(20), default='standalone')  # standalone, year_1, year_2

    # ✅ Future Stripe integration
    stripe_product_id = db.Column(db.String(100), nullable=True)
    stripe_price_id = db.Column(db.String(100), nullable=True)

    # ✅ Relationships for hierarchy
    children = db.relationship(
        'Course',
        backref=db.backref('parent', remote_side=[id]),
        lazy=True,
        cascade="all, delete-orphan"
    )

    lessons = db.relationship(
        "Lesson",
        backref="course",
        lazy=True,
        cascade="all, delete-orphan"   # ✅ Cascade delete lessons when course is deleted
    )

    progress_records = db.relationship(
        'UserCourseProgress',
        backref='course',
        lazy=True,
        cascade="all, delete-orphan"
    )

    quiz_retake_attempts = db.relationship(
        'QuizRetakeAttempt',
        backref='course',
        lazy=True,
        cascade="all, delete-orphan"
    )

    exams = db.relationship(
        'Exam',
        backref='course',
        lazy=True,
        cascade="all, delete-orphan"
    )

    course_agreements = db.relationship(
        'CourseAgreement',
        backref='course',
        lazy=True,
        cascade="all, delete-orphan"
    )

    course_access = db.relationship(
        'CourseAccess',
        backref='course',
        lazy=True,
        cascade="all, delete-orphan"
    )

    # ✅ Helper methods for hierarchy and access control
    def get_all_children(self):
        """Recursively get all child courses (years and sub-courses)"""
        result = []
        for child in self.children:
            result.append(child)
            result.extend(child.get_all_children())
        return result

    def is_parent_course(self):
        """Check if this is a parent course (has children)"""
        return len(self.children) > 0

    def is_free(self):
        """Check if course is free"""
        return self.price is None or self.price == 0

    def user_has_access(self, user_id):
        """Check if a user has access to this course"""
        # Check direct access
        direct_access = CourseAccess.query.filter_by(
            user_id=user_id,
            course_id=self.id
        ).first()
        if direct_access:
            return True

        # Check if user has access via parent course
        if self.parent_id:
            parent_access = CourseAccess.query.filter_by(
                user_id=user_id,
                course_id=self.parent_id
            ).first()
            if parent_access:
                return True

        return False


class Lesson(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    week = db.Column(db.Integer, nullable=False)
    title = db.Column(db.String(200), nullable=False)
    video_file = db.Column(db.String(300))
    ppt_file = db.Column(db.String(300))
    description = db.Column(db.Text, nullable=True)

    course_id = db.Column(db.Integer, db.ForeignKey("course.id"), nullable=False)

    quizzes = db.relationship(
        "Quiz",
        backref="lesson",
        lazy=True,
        cascade="all, delete-orphan"   # ✅ Cascade delete quizzes when lesson is deleted
    )

    quiz_attempts = db.relationship(
        "QuizAttempt",
        backref="lesson",
        lazy=True,
        cascade="all, delete-orphan"
    )

    quiz_retake_attempts = db.relationship(
        "QuizRetakeAttempt",
        backref="lesson",
        lazy=True,
        cascade="all, delete-orphan"
    )

    blocking_exams = db.relationship(
        'Exam',
        backref='trigger_lesson',
        lazy=True,
        cascade="all, delete-orphan",
        foreign_keys='Exam.trigger_lesson_id'
    )


class Exam(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'))
    title = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text)
    duration_minutes = db.Column(db.Integer, nullable=False, default=30)
    pass_mark = db.Column(db.Float, default=70.0)
    is_required = db.Column(db.Boolean, default=False)
    is_active = db.Column(db.Boolean, default=True)
    allow_retakes = db.Column(db.Boolean, default=True)
    trigger_lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'))
    required_to_complete_course = db.Column(db.Boolean, default=False)
    settings = db.Column(db.JSON, default=dict)

    questions = db.relationship(
        'ExamQuestion',
        backref='exam',
        lazy=True,
        cascade="all, delete-orphan",
        order_by='ExamQuestion.order_index'
    )

    attempts = db.relationship(
        'ExamAttempt',
        backref='exam',
        lazy=True,
        cascade="all, delete-orphan"
    )


class ExamQuestion(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    exam_id = db.Column(db.Integer, db.ForeignKey('exam.id'), nullable=False)
    question_type = db.Column(db.String(32), nullable=False, default='multiple_choice')
    text = db.Column(db.Text, nullable=False)
    options = db.Column(db.JSON, default=list)
    correct_answers = db.Column(db.JSON, default=list)
    points = db.Column(db.Float, nullable=False, default=1.0)
    is_required = db.Column(db.Boolean, default=True)
    order_index = db.Column(db.Integer, nullable=False, default=0)
    config = db.Column(db.JSON, default=dict)

    answers = db.relationship(
        'ExamAnswer',
        backref='question',
        lazy=True,
        cascade="all, delete-orphan"
    )


class ExamAttempt(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    exam_id = db.Column(db.Integer, db.ForeignKey('exam.id'), nullable=False)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'))
    start_time = db.Column(db.DateTime, default=utcnow)
    end_time = db.Column(db.DateTime)
    status = db.Column(db.String(20), default='in-progress')
    score = db.Column(db.Float, default=0.0)
    max_score = db.Column(db.Float, default=0.0)
    duration_seconds = db.Column(db.Integer)
    attempt_number = db.Column(db.Integer, default=1)
    autosave_payload = db.Column(db.JSON, default=dict)
    overall_feedback = db.Column(db.Text)
    passed = db.Column(db.Boolean)

    course = db.relationship('Course')

    answers = db.relationship(
        'ExamAnswer',
        backref='attempt',
        lazy=True,
        cascade="all, delete-orphan"
    )


class ExamAnswer(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    attempt_id = db.Column(db.Integer, db.ForeignKey('exam_attempt.id'), nullable=False)
    question_id = db.Column(db.Integer, db.ForeignKey('exam_question.id'), nullable=False)
    response_data = db.Column(db.JSON)
    is_correct = db.Column(db.Boolean)
    points_awarded = db.Column(db.Float, default=0.0)
    feedback = db.Column(db.Text)



class Quiz(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    question = db.Column(db.String(300), nullable=False)
    option_a = db.Column(db.String(200), nullable=False)
    option_b = db.Column(db.String(200), nullable=False)
    option_c = db.Column(db.String(200), nullable=False)
    correct_answer = db.Column(db.String(1), nullable=False)  # "A" / "B" / "C"

    lesson_id = db.Column(db.Integer, db.ForeignKey("lesson.id"), nullable=False)


class UserCourseProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)
    progress = db.Column(db.Integer, default=1)
    is_finished = db.Column(db.Boolean, default=False)  # ✅ Admin can mark as finished

    __table_args__ = (
        db.UniqueConstraint('user_id', 'course_id', name='uq_user_course'),
    )


# ✅ NEW: Course Access & Monetization Model
class CourseAccess(db.Model):
    """
    Tracks which users have access to which courses.
    Access can be granted via:
    - Direct purchase of a course
    - Purchase of a parent course (unlocks all children)
    - Admin grant
    - Free enrollment
    """
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)

    # Access metadata
    granted_at = db.Column(db.DateTime, default=utcnow)
    access_type = db.Column(db.String(20), default='free')  # free, purchased, admin_grant, parent_unlock, subscription
    is_locked = db.Column(db.Boolean, default=False)  # For subscription data retention
    progress = db.Column(db.Float, default=0.0)  # Track progress for data retention

    # Future Stripe payment tracking
    stripe_payment_intent_id = db.Column(db.String(100), nullable=True)
    amount_paid = db.Column(db.Numeric(10, 2), nullable=True)

    __table_args__ = (
        db.UniqueConstraint('user_id', 'course_id', name='uq_user_course_access'),
    )


class SubscriptionPlan(db.Model):
    """
    Defines monthly subscription plans that bundle multiple courses.
    """
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    price = db.Column(db.Numeric(10, 2), nullable=False)
    billing_interval = db.Column(db.String(20), default="monthly")  # monthly, yearly
    description = db.Column(db.Text)
    course_ids = db.Column(db.JSON)  # List of course IDs included in the plan
    is_active = db.Column(db.Boolean, default=True)
    stripe_price_id = db.Column(db.String(120))
    stripe_product_id = db.Column(db.String(120))
    trial_days = db.Column(db.Integer, default=7)  # 7-day free trial by default
    grace_period_days = db.Column(db.Integer, default=3)  # 3-day grace after failed payment
    created_at = db.Column(db.DateTime, default=utcnow)

    # Relationships
    subscriptions = db.relationship(
        'UserSubscription',
        backref='plan',
        lazy=True,
        cascade="all, delete-orphan"
    )

    def get_courses(self):
        """Get all Course objects included in this plan"""
        if not self.course_ids:
            return []
        return Course.query.filter(Course.id.in_(self.course_ids)).all()


class UserSubscription(db.Model):
    """
    Tracks individual user subscriptions to plans.
    """
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    plan_id = db.Column(db.Integer, db.ForeignKey('subscription_plan.id'), nullable=False)
    stripe_subscription_id = db.Column(db.String(120))
    stripe_customer_id = db.Column(db.String(120))
    status = db.Column(db.String(50))  # active, trialing, past_due, canceled, incomplete, incomplete_expired
    start_date = db.Column(db.DateTime, default=utcnow)
    end_date = db.Column(db.DateTime)
    current_period_start = db.Column(db.DateTime)
    current_period_end = db.Column(db.DateTime)
    cancel_at_period_end = db.Column(db.Boolean, default=False)
    canceled_at = db.Column(db.DateTime)
    trial_start = db.Column(db.DateTime)
    trial_end = db.Column(db.DateTime)

    # Payment tracking
    last_payment_date = db.Column(db.DateTime)
    last_payment_amount = db.Column(db.Numeric(10, 2))

    # Relationships
    user = db.relationship('User', backref=db.backref('subscriptions', lazy=True))

    def is_active(self):
        """Check if subscription is currently active (including trial and grace period)"""
        return self.status in ['active', 'trialing', 'past_due']

    def in_grace_period(self):
        """Check if subscription is in grace period after failed payment"""
        return self.status == 'past_due'


class StudentHubFile(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    stored_name = db.Column(db.String(300), nullable=False)
    original_name = db.Column(db.String(300), nullable=False)
    file_path = db.Column(db.String(400), nullable=False)
    file_size = db.Column(db.Integer, nullable=False)
    uploaded_at = db.Column(db.DateTime, default=utcnow)
    feedback_text = db.Column(db.Text)
    feedback_grade = db.Column(db.String(10))


class QuizAttempt(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    attempt_count = db.Column(db.Integer, default=0)
    last_score = db.Column(db.Integer, default=0)
    best_score = db.Column(db.Integer, default=0)
    total_questions = db.Column(db.Integer, default=0)
    correct_count = db.Column(db.Integer, default=0)
    wrong_count = db.Column(db.Integer, default=0)
    passed = db.Column(db.Boolean, default=False)
    detail_json = db.Column(db.Text)
    last_attempt_at = db.Column(db.DateTime, default=utcnow)

    __table_args__ = (
        db.UniqueConstraint('user_id', 'lesson_id', name='uq_user_lesson_attempt'),
    )


class QuizRetakeAttempt(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'))
    attempt_number = db.Column(db.Integer, nullable=False, default=1)
    score = db.Column(db.Integer, default=0)
    total_questions = db.Column(db.Integer, default=0)
    is_randomized = db.Column(db.Boolean, default=False)
    attempt_type = db.Column(db.String(20), default='lesson_list')
    is_complete = db.Column(db.Boolean, default=True)
    question_order_json = db.Column(db.Text)
    current_index = db.Column(db.Integer, default=0)
    answers_json = db.Column(db.Text)
    wrong_questions_json = db.Column(db.Text)
    created_at = db.Column(db.DateTime, default=utcnow)


class CourseAgreement(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)
    accepted_at = db.Column(db.DateTime, default=utcnow)

    __table_args__ = (
        db.UniqueConstraint('user_id', 'course_id', name='uq_course_agreement'),
    )


class Testimonial(db.Model):
    """Student testimonials/reviews for courses."""
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)
    name = db.Column(db.String(200), nullable=False)  # Default from user.full_name, editable
    rating = db.Column(db.Integer, nullable=False)  # 1-5 stars
    review = db.Column(db.Text, nullable=False)
    status = db.Column(db.String(20), default='pending')  # pending, approved, rejected
    created_at = db.Column(db.DateTime, default=utcnow)

    # Relationships
    user = db.relationship('User', backref='testimonials')
    course = db.relationship('Course', backref='testimonials')

    __table_args__ = (
        db.UniqueConstraint('user_id', 'course_id', name='uq_user_course_testimonial'),
    )


class SentEmail(db.Model):
    """Log of all emails sent by the system."""
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    subject = db.Column(db.String(200), nullable=False)
    sent_at = db.Column(db.DateTime, default=utcnow)
    status = db.Column(db.String(20), default='sent')  # sent, failed

    # Relationship
    user = db.relationship('User', backref='sent_emails')


class SiteSetting(db.Model):
    """Stores site-wide settings like Terms of Service, Privacy Policy, etc."""
    id = db.Column(db.Integer, primary_key=True)
    key = db.Column(db.String(100), unique=True, nullable=False)
    value = db.Column(db.Text, nullable=True)
    updated_at = db.Column(db.DateTime, default=utcnow, onupdate=utcnow)


def _student_hub_file_path(file_record: StudentHubFile) -> str:
    """Return absolute path for a stored student hub file."""
    return os.path.join(
        app.config['STUDENT_HUB_UPLOAD_FOLDER'],
        str(file_record.user_id),
        file_record.stored_name
    )


def _remove_student_hub_file_from_disk(file_record: StudentHubFile) -> None:
    """Best-effort removal of the stored file asset from disk."""
    try:
        os.remove(_student_hub_file_path(file_record))
    except FileNotFoundError:
        pass
    except OSError:
        pass


def _normalize_feedback_grade(raw_grade: str | None) -> str | None:
    """Return validated grade string or raise ValueError if invalid."""
    if not raw_grade:
        return None

    grade = raw_grade.strip()
    if not grade:
        return None

    numeric_match = re.fullmatch(r"\d{1,3}(?:\.\d+)?", grade)
    if numeric_match:
        value = float(grade)
        if 0 <= value <= 100:
            if value.is_integer():
                return str(int(value))
            return str(round(value, 2)).rstrip('0').rstrip('.')
        raise ValueError("Grade must be between 0 and 100.")

    grade_upper = grade.upper()
    allowed_letters = {"A", "B", "C", "D", "E", "F"}
    if grade_upper in allowed_letters:
        return grade_upper

    raise ValueError("Grade must be numeric (0-100) or one of A-F.")


def _fetch_course(course_id: int | None):
    return db.session.get(Course, course_id) if course_id else None


def _fetch_lesson(lesson_id: int | None):
    return db.session.get(Lesson, lesson_id) if lesson_id else None


def _fetch_quiz(quiz_id: int | None):
    return db.session.get(Quiz, quiz_id) if quiz_id else None


def _fetch_exam(exam_id: int | None):
    return db.session.get(Exam, exam_id) if exam_id else None


def _serialize_exam_question(question: 'ExamQuestion') -> dict:
    return {
        "id": question.id,
        "type": question.question_type,
        "text": question.text,
        "options": question.options or [],
        "correct_answers": question.correct_answers or [],
        "points": question.points or 0,
        "required": question.is_required,
        "config": question.config or {},
        "order_index": question.order_index or 0
    }


def _serialize_exam(exam: 'Exam') -> dict:
    return {
        "id": exam.id,
        "course_id": exam.course_id,
        "title": exam.title,
        "description": exam.description,
        "duration_minutes": exam.duration_minutes,
        "pass_mark": exam.pass_mark,
        "is_required": exam.is_required,
        "is_active": exam.is_active,
        "allow_retakes": exam.allow_retakes,
        "trigger_lesson_id": exam.trigger_lesson_id,
        "required_to_complete_course": exam.required_to_complete_course,
        "settings": exam.settings or {},
        "max_score": _exam_max_score(exam),
        "questions": [  
            _serialize_exam_question(question)
            for question in sorted(exam.questions, key=lambda q: (q.order_index or 0, q.id or 0))
        ]
    }


def _exam_max_score(exam: 'Exam') -> float:
    return float(sum(question.points or 0 for question in exam.questions)) or 0.0


def _exam_grading_mode(exam: 'Exam') -> str:
    settings = exam.settings or {}
    mode = (settings.get('grading_mode') or 'automatic').lower()
    if mode not in {'automatic', 'manual'}:
        return 'automatic'
    return mode


def _grade_objective_response(question: 'ExamQuestion', response_data, auto_pass_subjective: bool = False) -> tuple[bool | None, float, dict]:
    """Return (is_correct, points_awarded, normalized_response).

    For non-objective question types we return (None, 0, normalized).
    If auto_pass_subjective is True, subjective questions are awarded full points.
    """
    question_type = (question.question_type or '').lower()

    if question_type in {'multiple_choice', 'mcq', 'radio'}:
        selected = None
        if isinstance(response_data, str):
            selected = response_data
        elif isinstance(response_data, dict):
            selected = response_data.get('selected') or response_data.get('value')

        correct_options = question.correct_answers or []
        normalized = {
            'selected': selected
        }

        if not correct_options:
            if auto_pass_subjective:
                points = float(question.points or 0)
                return True, points, normalized
            return None, 0.0, normalized

        is_correct = selected in correct_options
        points = float(question.points or 0)
        return is_correct, (points if is_correct else 0.0), normalized

    if question_type in {'checkbox', 'multi_select'}:
        if isinstance(response_data, dict):
            selected = response_data.get('selected') or response_data.get('values') or []
        elif isinstance(response_data, (list, tuple, set)):
            selected = list(response_data)
        else:
            selected = [response_data] if response_data else []

        selected_set = {str(item) for item in selected if item is not None}
        correct_options = {str(item) for item in (question.correct_answers or [])}
        normalized = {
            'selected': sorted(selected_set)
        }

        if not correct_options:
            return None, 0.0, normalized

        is_correct = selected_set == correct_options
        points = float(question.points or 0)
        return is_correct, (points if is_correct else 0.0), normalized

    if question_type == 'short_answer':
        text = ''
        if isinstance(response_data, str):
            text = response_data.strip()
        elif isinstance(response_data, dict):
            text = (response_data.get('text') or '').strip()
        normalized = {
            'text': text
        }

        expected_answers = [
            str(answer).strip().lower()
            for answer in (question.correct_answers or [])
            if isinstance(answer, str)
        ]

        if expected_answers and text:
            is_correct = text.lower() in expected_answers
            points = float(question.points or 0)
            return is_correct, (points if is_correct else 0.0), normalized

        if auto_pass_subjective:
            points = float(question.points or 0)
            return True, points, normalized
        return None, 0.0, normalized

    if question_type in {'essay', 'long_answer'}:
        text = ''
        if isinstance(response_data, str):
            text = response_data
        elif isinstance(response_data, dict):
            text = response_data.get('text') or ''
        normalized = {
            'text': text
        }
        if auto_pass_subjective:
            points = float(question.points or 0)
            return True, points, normalized
        return None, 0.0, normalized

    return None, 0.0, {'value': response_data}


def _compute_attempt_score(attempt: 'ExamAttempt') -> None:
    """Aggregate scoring fields for the attempt and set status if auto-graded."""
    total_awarded = 0.0
    requires_manual = False
    for answer in attempt.answers:
        total_awarded += float(answer.points_awarded or 0)
        if answer.is_correct is None:
            requires_manual = True

    attempt.score = round(total_awarded, 2)
    attempt.max_score = _exam_max_score(attempt.exam)
    attempt.passed = None

    if attempt.max_score > 0:
        percent = (attempt.score / attempt.max_score) * 100
        pass_mark = attempt.exam.pass_mark or 0
        if not requires_manual:
            attempt.passed = percent >= pass_mark
            attempt.status = 'graded'
        else:
            attempt.status = 'submitted'
    else:
        attempt.status = 'submitted'


def _finish_attempt(attempt: 'ExamAttempt', submitted_at: datetime | None = None) -> None:
    attempt.start_time = _ensure_utc(attempt.start_time) or utcnow()
    attempt.end_time = _ensure_utc(submitted_at) or utcnow()
    attempt.duration_seconds = int((attempt.end_time - attempt.start_time).total_seconds())
    if attempt.duration_seconds < 0:
        attempt.duration_seconds = 0
    _compute_attempt_score(attempt)


def _time_remaining_seconds(attempt: 'ExamAttempt') -> int:
    attempt.start_time = _ensure_utc(attempt.start_time) or utcnow()
    duration = int((attempt.exam.duration_minutes or 0) * 60)
    elapsed = int((utcnow() - attempt.start_time).total_seconds())
    remaining = duration - elapsed
    return remaining if remaining > 0 else 0


def _user_has_passed_exam(user_id: int, exam_id: int) -> bool:
    passed_attempt = ExamAttempt.query.filter(
        ExamAttempt.user_id == user_id,
        ExamAttempt.exam_id == exam_id,
        ExamAttempt.status.in_(['passed', 'graded'])
    ).order_by(ExamAttempt.id.desc()).first()
    return bool(passed_attempt)


def _active_attempt_for_user(user_id: int, exam_id: int) -> 'ExamAttempt | None':
    attempt = ExamAttempt.query.filter(
        ExamAttempt.user_id == user_id,
        ExamAttempt.exam_id == exam_id,
        ExamAttempt.status == 'in-progress'
    ).order_by(ExamAttempt.start_time.desc()).first()
    if attempt:
        original_start = attempt.start_time
        attempt.start_time = _ensure_utc(attempt.start_time) or utcnow()
        if original_start != attempt.start_time:
            db.session.commit()
        remaining = _time_remaining_seconds(attempt)
        if remaining <= 0:
            _finish_attempt(attempt)
            db.session.commit()
            return None
    return attempt


def _start_exam_attempt(user: User, exam: 'Exam') -> 'ExamAttempt':
    existing_attempt = _active_attempt_for_user(user.id, exam.id)
    if existing_attempt:
        return existing_attempt

    attempt_count = ExamAttempt.query.filter_by(user_id=user.id, exam_id=exam.id).count()
    new_attempt = ExamAttempt(
        user_id=user.id,
        exam_id=exam.id,
        course_id=exam.course_id,
        attempt_number=attempt_count + 1,
        status='in-progress'
    )
    db.session.add(new_attempt)
    db.session.flush()
    new_attempt.max_score = _exam_max_score(exam)
    new_attempt.start_time = _ensure_utc(new_attempt.start_time) or utcnow()
    return new_attempt


def _summarize_exam_attempt(attempt: 'ExamAttempt') -> dict:
    exam = attempt.exam
    summary = {
        'attempt_id': attempt.id,
        'exam_id': attempt.exam_id,
        'user_id': attempt.user_id,
        'status': attempt.status,
        'score': attempt.score,
        'max_score': attempt.max_score,
        'percentage': round((attempt.score / attempt.max_score) * 100, 2) if attempt.max_score else None,
        'start_time': attempt.start_time.isoformat() if attempt.start_time else None,
        'end_time': attempt.end_time.isoformat() if attempt.end_time else None,
        'duration_seconds': attempt.duration_seconds,
        'overall_feedback': attempt.overall_feedback,
        'passed': attempt.passed,
        'requires_manual_grading': any(answer.is_correct is None for answer in attempt.answers),
        'answers': []
    }

    question_lookup = {question.id: question for question in (exam.questions if exam else [])}
    for answer in attempt.answers:
        question = question_lookup.get(answer.question_id)
        summary['answers'].append({
            'id': answer.id,
            'question_id': answer.question_id,
            'question_text': question.text if question else None,
            'question_type': question.question_type if question else None,
            'response_data': answer.response_data,
            'is_correct': answer.is_correct,
            'points_awarded': answer.points_awarded,
            'points_possible': question.points if question else None,
            'feedback': answer.feedback,
            'correct_answers': question.correct_answers if question else None,
            'options': question.options if question else None
        })

    return summary


def _exam_statistics(exam_id: int) -> dict:
    aggregates = db.session.query(
        func.count(ExamAttempt.id),
        func.avg(ExamAttempt.score),
        func.max(ExamAttempt.score),
        func.min(ExamAttempt.score),
        func.avg(ExamAttempt.duration_seconds),
        func.sum(case((ExamAttempt.passed == True, 1), else_=0))  # noqa: E712
    ).filter(
        ExamAttempt.exam_id == exam_id,
        ExamAttempt.status != 'in-progress'
    ).one()

    attempt_count = aggregates[0] or 0
    avg_score = float(aggregates[1]) if aggregates[1] is not None else None
    max_score = float(aggregates[2]) if aggregates[2] is not None else None
    min_score = float(aggregates[3]) if aggregates[3] is not None else None
    avg_duration = float(aggregates[4]) if aggregates[4] is not None else None
    passed_total = int(aggregates[5] or 0)

    pass_rate = None
    if attempt_count:
        pass_rate = round((passed_total / attempt_count) * 100, 2)

    question_rows = db.session.query(
        ExamAnswer.question_id,
        func.count(ExamAnswer.id).label('total'),
        func.sum(case((ExamAnswer.is_correct == True, 1), else_=0)).label('correct')  # noqa: E712
    ).join(ExamQuestion).filter(
        ExamQuestion.exam_id == exam_id
    ).group_by(ExamAnswer.question_id).all()

    most_missed = None
    for question_id, total, correct in question_rows:
        total = total or 0
        correct = correct or 0
        missed = total - correct
        if total == 0:
            continue
        ratio = missed / total
        if not most_missed or ratio > most_missed['ratio']:
            question = db.session.get(ExamQuestion, question_id)
            most_missed = {
                'question_id': question_id,
                'question_text': question.text if question else None,
                'missed': missed,
                'total': total,
                'ratio': ratio
            }

    return {
        'attempt_count': attempt_count,
        'average_score': round(avg_score, 2) if avg_score is not None else None,
        'highest_score': max_score,
        'lowest_score': min_score,
        'average_duration': round(avg_duration, 2) if avg_duration is not None else None,
        'pass_rate': pass_rate,
        'passed_total': passed_total,
        'most_missed': most_missed
    }


def _admin_attempt_payload(attempt: 'ExamAttempt') -> dict:
    exam = attempt.exam
    user = attempt.user
    payload = _summarize_exam_attempt(attempt)
    payload.update({
        'user': {
            'id': user.id if user else None,
            'name': user.full_name or user.username if user else 'Unknown',
            'email': user.email if user else None
        },
        'status': attempt.status,
        'attempt_number': attempt.attempt_number,
        'autosave_payload': attempt.autosave_payload,
        'course_id': attempt.course_id
    })

    question_lookup = {question.id: question for question in (exam.questions if exam else [])}
    for answer in payload['answers']:
        question = question_lookup.get(answer['question_id'])
        answer['question_order'] = question.order_index if question else None

    payload['answers'].sort(key=lambda item: item.get('question_order') or 0)
    return payload

def _extract_text_from_upload(upload):
    """Return textual content parsed from a lesson upload, if possible."""
    if not upload or not getattr(upload, 'filename', None):
        return ""

    filename = upload.filename or ""
    extension = os.path.splitext(filename)[1].lower()

    try:
        binary_data = upload.read()
    finally:
        try:
            upload.stream.seek(0)
        except Exception:
            try:
                upload.seek(0)
            except Exception:
                pass

    if not binary_data:
        return ""

    try:
        if extension in {'.ppt', '.pptx'}:
            prs = Presentation(io.BytesIO(binary_data))
            slide_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        slide_text.append(shape.text)
            return "\n".join(slide_text)

        if extension == '.pdf':
            with fitz.open(stream=binary_data, filetype='pdf') as doc:
                page_text = [page.get_text('text') for page in doc]
            return "\n".join(page_text)

        if extension in {'.txt', '.md'}:
            return binary_data.decode('utf-8', errors='ignore')

        # Fallback: attempt to decode as UTF-8 text
        return binary_data.decode('utf-8', errors='ignore')
    except Exception:
        return ""


def _build_quiz_from_text(raw_text: str, requested: int = 3):
    """Generate simple multiple-choice questions from lesson text."""
    cleaned = re.sub(r'\s+', ' ', raw_text or '').strip()
    if not cleaned:
        return []

    sentences = [
        sentence.strip()
        for sentence in re.split(r'(?<=[.!?])\s+', cleaned)
        if len(sentence.strip()) >= 40
    ]

    if len(sentences) < 3:
        # fall back to line breaks if too few sentences
        sentences = [
            line.strip()
            for line in re.split(r'[\r\n]+', cleaned)
            if len(line.strip()) >= 30
        ]

    # De-duplicate while preserving order
    unique_sentences = []
    seen = set()
    for sentence in sentences:
        key = sentence.lower()
        if key not in seen:
            unique_sentences.append(sentence)
            seen.add(key)

    if not unique_sentences:
        return []

    pool = unique_sentences[:]
    random.shuffle(pool)
    num_questions = max(1, min(requested, len(pool)))
    questions = []
    placeholder_pool = [
        "This statement is not covered in the lesson.",
        "An unrelated fact introduced for comparison.",
        "A detail from a different topic entirely."
    ]

    for idx in range(num_questions):
        correct_statement = pool[idx % len(pool)]
        distractor_source = [s for s in unique_sentences if s != correct_statement]

        distractors = []
        if distractor_source:
            sample_size = min(2, len(distractor_source))
            distractors = random.sample(distractor_source, sample_size)

        while len(distractors) < 2:
            choice = random.choice(placeholder_pool)
            if choice not in distractors and choice.lower() != correct_statement.lower():
                distractors.append(choice)

        option_entries = [
            {"text": correct_statement, "is_correct": True},
            {"text": distractors[0], "is_correct": False},
            {"text": distractors[1], "is_correct": False}
        ]
        random.shuffle(option_entries)

        option_map = {}
        correct_letter = 'A'
        letters = ['A', 'B', 'C']
        for pos, entry in enumerate(option_entries):
            trimmed_text = entry["text"].strip()
            if len(trimmed_text) > 220:
                trimmed_text = trimmed_text[:217].rstrip() + '...'
            option_map[letters[pos]] = trimmed_text
            if entry["is_correct"]:
                correct_letter = letters[pos]

        question_text = f"Identify the statement from the lesson (item {idx + 1})."
        questions.append({
            "question": question_text,
            "option_a": option_map['A'],
            "option_b": option_map['B'],
            "option_c": option_map['C'],
            "correct_answer": correct_letter
        })

    return questions
    # -------------------- Q&A Chat Models -------------------- #
class Question(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id', ondelete='CASCADE'), nullable=False)
    title = db.Column(db.String(200), nullable=False)
    is_public = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.DateTime, default=db.func.current_timestamp())

    user = db.relationship("User", backref=db.backref("questions", cascade="all, delete-orphan"))
    messages = db.relationship("Message", backref="question", lazy=True, cascade="all, delete-orphan")



class Message(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    question_id = db.Column(db.Integer, db.ForeignKey('question.id'), nullable=False)
    sender = db.Column(db.String(50))   # "user" or "admin"
    body = db.Column(db.Text, nullable=False)
    created_at = db.Column(db.DateTime, default=db.func.current_timestamp())


# -------------------- Subscription Helper Functions -------------------- #
def grant_course_access(user_id, course, *, access_type='purchased', amount_paid=None, payment_intent_id=None):
    """Grant access to a course and all published children."""

    def _grant(target, record_type, record_amount, record_payment):
        existing_access = CourseAccess.query.filter_by(
            user_id=user_id,
            course_id=target.id
        ).first()

        if existing_access:
            existing_access.is_locked = False
            if record_type:
                existing_access.access_type = record_type
            if record_payment:
                existing_access.stripe_payment_intent_id = record_payment
            if record_amount is not None:
                existing_access.amount_paid = record_amount
        else:
            db.session.add(CourseAccess(
                user_id=user_id,
                course_id=target.id,
                access_type=record_type or 'purchased',
                is_locked=False,
                amount_paid=record_amount,
                stripe_payment_intent_id=record_payment
            ))

        for child in target.children:
            if child.is_published:
                _grant(child, 'parent_unlock', None, None)

    _grant(course, access_type, amount_paid, payment_intent_id)


def grant_subscription_access(user_id, course_ids):
    """
    Grant subscription access to multiple courses for a user.
    Creates CourseAccess records with subscription type.
    """
    for course_id in course_ids:
        course = Course.query.get(course_id)
        if course:
            grant_course_access(user_id, course, access_type='subscription')


def revoke_subscription_access(user_id, course_ids, keep_progress=True):
    """
    Revoke subscription access to courses.
    If keep_progress is True, locks access but retains data.
    If False, removes access entirely.
    """
    for course_id in course_ids:
        access = CourseAccess.query.filter_by(
            user_id=user_id,
            course_id=course_id,
            access_type='subscription'
        ).first()

        if access:
            if keep_progress:
                access.is_locked = True
            else:
                db.session.delete(access)


def unlock_subscription_courses(user_id, course_ids):
    """
    Unlock previously locked subscription courses (e.g., after payment recovery).
    """
    for course_id in course_ids:
        access = CourseAccess.query.filter_by(
            user_id=user_id,
            course_id=course_id
        ).first()

        if access:
            access.is_locked = False


def user_has_active_subscription(user_id):
    """
    Check if user has any active subscription.
    """
    active_sub = UserSubscription.query.filter_by(
        user_id=user_id
    ).filter(
        UserSubscription.status.in_(['active', 'trialing', 'past_due'])
    ).first()

    return active_sub is not None


def get_user_active_subscription(user_id):
    """
    Get user's active subscription (if any).
    Returns UserSubscription object or None.
    """
    return UserSubscription.query.filter_by(
        user_id=user_id
    ).filter(
        UserSubscription.status.in_(['active', 'trialing', 'past_due'])
    ).first()


# -------------------- Session Decorators -------------------- #
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user' not in session:
            flash("You must log in first!")
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

def user_only(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if session.get('role') != 'user':
            flash("Access denied!")
            return redirect(url_for('index'))
        return f(*args, **kwargs)
    return decorated_function

def admin_only(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if session.get('role') != 'admin':
            flash("Access denied!")
            return redirect(url_for('index'))
        return f(*args, **kwargs)
    return decorated_function

# -------------------- Routes -------------------- #
@app.route('/')
def index():
    public_qs = Question.query.filter_by(is_public=True).all()

    # ✅ NEW: Only show courses marked to display on homepage
    courses = Course.query.filter_by(
        show_on_homepage=True,
        is_published=True
    ).order_by(Course.order_index, Course.name).all()

    # fetch approved testimonials for homepage
    testimonials = Testimonial.query.filter_by(status='approved').order_by(Testimonial.created_at.desc()).all()

    # Get current user if logged in
    current_user = None
    if 'user' in session:
        current_user = User.query.filter_by(username=session['user']).first()

    return render_template('index.html', public_qs=public_qs, courses=courses, testimonials=testimonials, current_user=current_user)


# -------------------- Terms of Service -------------------- #
@app.route('/terms')
def terms():
    """Display Terms of Service page."""
    terms_setting = SiteSetting.query.filter_by(key='terms_of_service').first()
    terms_content = terms_setting.value if terms_setting else "Terms of Service not yet configured."
    return render_template('terms.html', terms_content=terms_content)


# -------------------- Registration -------------------- #
# -------------------- Registration -------------------- #
@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        # ✅ normalize username to lowercase
        username = request.form['username'].strip().lower()
        raw_password = request.form['password']
        password = bcrypt.generate_password_hash(raw_password).decode('utf-8')
        full_name = request.form.get('full_name')
        age_raw = request.form.get('age', '').strip()
        if age_raw:
            try:
                age = int(age_raw)
            except ValueError:
                flash("Please enter a valid numeric age or leave the field blank.")
                return redirect(url_for('register'))
        else:
            age = None
        phone_number = request.form.get('phone_number')
        email = request.form.get('email')

        # Validate Terms of Service acceptance
        accept_terms = request.form.get('accept_terms')
        if not accept_terms:
            flash("You must accept the Terms of Service to register.")
            return redirect(url_for('register'))

        # Check if username or email already exists
        if User.query.filter((User.username == username) | (User.email == email)).first():
            flash("Account already exists!")
            return redirect(url_for('register'))

        # Create and save new user
        new_user = User(
            username=username,  # always lowercase
            password=password,
            full_name=full_name,
            age=age,
            phone_number=phone_number,
            email=email,
            role="user",
            accepted_terms=True,
            terms_accepted_at=datetime.now(timezone.utc)
        )
        db.session.add(new_user)
        db.session.commit()

        # Send welcome email
        try:
            from email_utils import send_welcome_email
            send_welcome_email(new_user)
        except Exception as e:
            # Log error but don't block registration
            print(f"Failed to send welcome email: {e}")

        flash("Registration successful! You can now log in.")
        return redirect(url_for('login'))

    return render_template('registration.html')


# -------------------- Login -------------------- #
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'GET':
        return render_template("log_in.html")

    # POST login logic
    username_or_email = request.form['username'].strip().lower()
    password = request.form['password']

    # Check if login is email or username
    user = User.query.filter(
        or_(User.username == username_or_email, User.email == username_or_email)
    ).first()

    if user and bcrypt.check_password_hash(user.password, password):
        user.last_login = datetime.now(timezone.utc)
        db.session.commit()
        session['user'] = user.username
        session['role'] = user.role

        if user.role == "admin":
            return redirect(url_for('admin_success'))
        elif user.role in ["user", "paid"]:
            return redirect(url_for('courses_dashboard'))
        else:
            return redirect(url_for('index'))
    else:
        flash("Invalid Username/Email or Password")
        return redirect(url_for('login'))
# -------------------- Logout -------------------- #
@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))


# -------------------- Password Reset -------------------- #
@app.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    """Request password reset - sends email with reset link"""
    if request.method == 'GET':
        return render_template('forgot_password.html')

    # POST: Process password reset request
    email = request.form.get('email', '').strip().lower()

    if not email:
        flash('Please enter your email address', 'error')
        return redirect(url_for('forgot_password'))

    # Find user by email
    user = User.query.filter_by(email=email).first()

    if user:
        # Generate secure token (expires in 1 hour)
        token = serializer.dumps(email, salt='password-reset-salt')

        # Store token and expiry in database (additional security layer)
        user.reset_token = token
        user.reset_token_expiry = datetime.now() + timedelta(hours=1)
        db.session.commit()

        # Send password reset email using new email utility
        try:
            from email_utils import send_password_reset_email
            send_password_reset_email(user, token)
            flash('Password reset instructions have been sent to your email address. Please check your inbox.', 'success')

        except Exception as e:
            print(f"Error sending email: {e}")
            flash('An error occurred while sending the reset email. Please try again later.', 'error')
            return redirect(url_for('forgot_password'))
    else:
        # Don't reveal if email exists (security best practice)
        # Show same message whether email exists or not
        flash('If an account with that email exists, password reset instructions have been sent.', 'info')

    return redirect(url_for('login'))


@app.route('/reset-password/<token>', methods=['GET', 'POST'])
def reset_password(token):
    """Reset password using token from email"""

    # Verify token is valid and not expired (itsdangerous check)
    try:
        email = serializer.loads(token, salt='password-reset-salt', max_age=3600)  # 1 hour expiry
    except SignatureExpired:
        flash('This password reset link has expired. Please request a new one.', 'error')
        return redirect(url_for('forgot_password'))
    except BadSignature:
        flash('Invalid password reset link. Please request a new one.', 'error')
        return redirect(url_for('forgot_password'))

    # Find user and verify database token matches (prevent token reuse)
    user = User.query.filter_by(email=email).first()

    if not user:
        flash('Invalid password reset link.', 'error')
        return redirect(url_for('forgot_password'))

    # Check if token matches and hasn't expired in database
    if user.reset_token != token:
        flash('This password reset link has already been used. Please request a new one.', 'error')
        return redirect(url_for('forgot_password'))

    if user.reset_token_expiry and user.reset_token_expiry < datetime.now():
        flash('This password reset link has expired. Please request a new one.', 'error')
        return redirect(url_for('forgot_password'))

    # GET: Show password reset form
    if request.method == 'GET':
        return render_template('reset_password.html', token=token, email=email)

    # POST: Process new password
    new_password = request.form.get('password', '').strip()
    confirm_password = request.form.get('confirm_password', '').strip()

    # Validation
    if not new_password or not confirm_password:
        flash('Please fill in all fields', 'error')
        return render_template('reset_password.html', token=token, email=email)

    if new_password != confirm_password:
        flash('Passwords do not match', 'error')
        return render_template('reset_password.html', token=token, email=email)

    if len(new_password) < 6:
        flash('Password must be at least 6 characters long', 'error')
        return render_template('reset_password.html', token=token, email=email)

    # Hash new password
    hashed_password = bcrypt.generate_password_hash(new_password).decode('utf-8')

    # Update user password and invalidate token
    user.password = hashed_password
    user.reset_token = None  # Invalidate token to prevent reuse
    user.reset_token_expiry = None
    db.session.commit()

    flash('Your password has been successfully reset. You can now log in with your new password.', 'success')
    return redirect(url_for('login'))


# -------------------- Success Pages -------------------- #
@app.route('/user_success')
@login_required
@user_only
def user_success():
    user = session['user']
    return render_template('home.html', user=user)


@app.route('/admin_success')
@login_required
@admin_only
def admin_success():
    return redirect(url_for('admin_dashboard'))

# -------------------- Course Management -------------------- #
@app.route('/admin/courses')
@login_required
@admin_only
def manage_courses():
    courses = Course.query.all()
    lessons = Lesson.query.order_by(Lesson.course_id, Lesson.week).all()
    return render_template("admin_courses.html", courses=courses, lessons=lessons)

@app.route('/admin/lesson/<int:lesson_id>/edit', methods=['GET', 'POST'])
@login_required
@admin_only
def edit_lesson(lesson_id):
    lesson = Lesson.query.get_or_404(lesson_id)

    if request.method == "POST":
        lesson.title = request.form.get("title")
        lesson.week = int(request.form.get("week"))
        lesson.description = request.form.get("description")
        # handle replacing files if uploaded again
        video_file = request.files.get("video_file")
        ppt_file = request.files.get("ppt_file")
        if video_file:
            folder = os.path.join("static", "courses", lesson.course.name, f"year{lesson.course.year}")
            os.makedirs(folder, exist_ok=True)
            path = os.path.join(folder, video_file.filename)
            video_file.save(path)
            lesson.video_file = os.path.join("courses", lesson.course.name, f"year{lesson.course.year}", video_file.filename)
        if ppt_file:
            folder = os.path.join("static", "courses", lesson.course.name, f"year{lesson.course.year}")
            os.makedirs(folder, exist_ok=True)
            path = os.path.join(folder, ppt_file.filename)
            ppt_file.save(path)
            lesson.ppt_file = os.path.join("courses", lesson.course.name, f"year{lesson.course.year}", ppt_file.filename)

        db.session.commit()
        flash("✅ Lesson updated successfully")
        return redirect(url_for('manage_courses'))

    return render_template("edit_lesson.html", lesson=lesson)


@app.route('/admin/lesson/<int:lesson_id>/quiz/update', methods=['POST'])
@login_required
@admin_only
def update_lesson_quiz(lesson_id):
    lesson = Lesson.query.get_or_404(lesson_id)
    raw_data = request.form.get('quiz_data', '').strip()

    Quiz.query.filter_by(lesson_id=lesson.id).delete()

    if raw_data:
        for line in raw_data.splitlines():
            parts = [p.strip() for p in line.split(',')]
            if len(parts) < 5:
                continue
            q_text, opt_a, opt_b, opt_c, correct_num = parts[:5]
            correct_letter = {"1": "A", "2": "B", "3": "C"}.get(correct_num, "A")
            quiz = Quiz(
                question=q_text,
                option_a=opt_a,
                option_b=opt_b,
                option_c=opt_c,
                correct_answer=correct_letter,
                lesson_id=lesson.id
            )
            db.session.add(quiz)

    db.session.commit()
    flash("✅ Quiz updated for this lesson")
    return redirect(url_for('edit_lesson', lesson_id=lesson.id))


@app.route('/admin/courses/upload', methods=['POST'])
@login_required
@admin_only
def upload_course_file():
    course_name = request.form.get("course")
    year = int(request.form.get("year"))
    week = int(request.form.get("week"))
    title = request.form.get("title")
    description = request.form.get("description")

    video_file = request.files.get("video_file")
    ppt_file = request.files.get("ppt_file")
    quiz_data = request.form.get("quiz_data")

    if not course_name or not year or not week or not title:
        flash("⚠️ Please fill all required fields.")
        return redirect(url_for("manage_courses"))


    # ✅ Find or create course
    course = Course.query.filter_by(name=course_name, year=year).first()
    if not course:
        course = Course(name=course_name, year=year)
        db.session.add(course)
        db.session.commit()

    # ✅ Find or create lesson
    lesson = Lesson.query.filter_by(course_id=course.id, week=week).first()
    if not lesson:
        lesson = Lesson(week=week, title=title, description=description, course_id=course.id)
    else:
        lesson.title = title or lesson.title
        lesson.description = description if description is not None else lesson.description

    # ✅ Save uploaded files
    folder = os.path.join("static", "courses", course_name, f"year{year}")
    os.makedirs(folder, exist_ok=True)

    if video_file:
        video_path = os.path.join(folder, video_file.filename)
        video_file.save(video_path)
        lesson.video_file = os.path.join("courses", course_name, f"year{year}", video_file.filename)

    if ppt_file:
        ppt_path = os.path.join(folder, ppt_file.filename)
        ppt_file.save(ppt_path)
        lesson.ppt_file = os.path.join("courses", course_name, f"year{year}", ppt_file.filename)

    db.session.add(lesson)
    db.session.commit()

    # ✅ Handle quizzes (CSV-style input)
    if quiz_data:
        for line in quiz_data.strip().splitlines():
            parts = [p.strip() for p in line.split(",")]
            if len(parts) == 5:  # Question, A, B, C, CorrectNum
                q_text, opt_a, opt_b, opt_c, correct_num = parts
                correct_letter = {"1": "A", "2": "B", "3": "C"}.get(correct_num, "A")
                quiz = Quiz(
                    question=q_text,
                    option_a=opt_a,
                    option_b=opt_b,
                    option_c=opt_c,
                    correct_answer=correct_letter,
                    lesson_id=lesson.id
                )
                db.session.add(quiz)

        db.session.commit()

    flash(f"✅ Lesson {title} uploaded for {course_name} (Year {year}, Week {week})")
    return redirect(url_for("manage_courses"))


@app.route('/admin/courses/generate', methods=['POST'])
@login_required
@admin_only
def generate_quiz_ajax():
    ppt_file = request.files.get('ppt_file') or request.files.get('slides')
    num_questions_raw = request.form.get('num_questions', '3')

    try:
        num_questions = max(1, min(20, int(num_questions_raw)))
    except (TypeError, ValueError):
        num_questions = 3

    if not ppt_file:
        return jsonify({
            "error": "Please attach slides (PPTX/PDF/TXT) before generating quiz questions."
        }), 400

    raw_text = _extract_text_from_upload(ppt_file)
    if not raw_text:
        return jsonify({
            "error": "We couldn't read any text content from the uploaded file. Please ensure it is a PPTX or PDF with selectable text."
        }), 400

    questions = _build_quiz_from_text(raw_text, requested=num_questions)
    if not questions:
        return jsonify({
            "error": "Not enough readable content to draft quiz questions. Try adding more detailed slides or enter questions manually."}
        ), 400

    return jsonify({
        "questions": questions
    })


# -------------------- Admin Dashboard -------------------- #
# -------------------- Admin Dashboard -------------------- #
@app.route('/admin')
@app.route('/admin/dashboard')
@login_required
@admin_only
def admin_dashboard():
    return render_template("admin_dashboard.html")


@app.route('/admin/student-files')
@login_required
@admin_only
def admin_student_files():
    courses = Course.query.order_by(Course.year, Course.name).all()

    student_filter = request.args.get('student_id', default=None, type=int)
    course_filter = request.args.get('course_name', default='', type=str)
    search_query = request.args.get('search', default='', type=str)

    files_query = StudentHubFile.query.join(User).order_by(StudentHubFile.uploaded_at.desc())

    if student_filter:
        files_query = files_query.filter(StudentHubFile.user_id == student_filter)

    if course_filter:
        like_filter = f"%{course_filter.lower()}%"
        files_query = files_query.filter(
            func.lower(func.coalesce(User.courses, '')).like(like_filter)
        )

    if search_query:
        like_search = f"%{search_query.lower()}%"
        files_query = files_query.filter(
            or_(
                func.lower(func.coalesce(User.full_name, '')).like(like_search),
                func.lower(User.username).like(like_search),
                func.lower(User.email).like(like_search)
            )
        )

    student_hub_files = files_query.all()
    students_with_files = sorted(
        User.query.join(StudentHubFile).distinct().all(),
        key=lambda u: (u.full_name or u.username).lower()
    )
    course_names = sorted({course.name for course in courses})

    return render_template(
        "admin_student_files.html",
        student_hub_files=student_hub_files,
        students_with_files=students_with_files,
        course_names=course_names,
        selected_student_id=student_filter,
        selected_course_name=course_filter,
        selected_search=search_query
    )


@app.route('/admin/exams')
@login_required
@admin_only
def admin_exams():
    exams = Exam.query.order_by(Exam.is_active.desc(), Exam.title.asc()).all()

    stats_lookup = {
        row.exam_id: row
        for row in db.session.query(
            ExamAttempt.exam_id,
            func.count(ExamAttempt.id).label('attempts'),
            func.avg(ExamAttempt.score).label('avg_score'),
            func.sum(case((ExamAttempt.passed == True, 1), else_=0)).label('passed_total')  # noqa: E712
        ).filter(
            ExamAttempt.status != 'in-progress'
        ).group_by(ExamAttempt.exam_id).all()
    }

    exams_payload = []
    for exam in exams:
        stats = stats_lookup.get(exam.id)
        attempts = int(stats.attempts) if stats else 0
        avg_score = float(stats.avg_score) if stats and stats.avg_score is not None else None
        passed_total = int(stats.passed_total) if stats else 0

        exams_payload.append({
            'exam': exam,
            'course': exam.course,
            'question_count': len(exam.questions),
            'attempts': attempts,
            'average_score': round(avg_score, 2) if avg_score is not None else None,
            'pass_rate': round((passed_total / attempts) * 100, 1) if attempts else None
        })

    return render_template(
        'admin_exams.html',
        exams_payload=exams_payload
    )


def _builder_context(exam: 'Exam | None' = None) -> dict:
    courses = Course.query.order_by(Course.year, Course.name).all()
    course_payload = []
    for course in courses:
        course_payload.append({
            'id': course.id,
            'name': course.name,
            'year': course.year,
            'lessons': [
                {
                    'id': lesson.id,
                    'title': lesson.title,
                    'week': lesson.week
                }
                for lesson in sorted(course.lessons, key=lambda l: l.week)
            ]
        })

    return {
        'exam_data': _serialize_exam(exam) if exam else None,
        'courses': course_payload
    }


@app.route('/admin/exams/new')
@login_required
@admin_only
def admin_exam_new():
    context = _builder_context()
    return render_template(
        'admin_exam_edit.html',
        exam=context['exam_data'],
        courses_json=context['courses'],
        is_new=True
    )


@app.route('/admin/exams/<int:exam_id>/edit')
@login_required
@admin_only
def admin_exam_edit(exam_id):
    exam = Exam.query.get_or_404(exam_id)
    context = _builder_context(exam)
    return render_template(
        'admin_exam_edit.html',
        exam=context['exam_data'],
        courses_json=context['courses'],
        is_new=False,
        linked_course_id=exam.course_id
    )


def _apply_exam_updates(exam: Exam, exam_data: dict, questions_data: list[dict]):
    exam.title = exam_data.get('title', exam.title)
    exam.description = exam_data.get('description')
    exam.duration_minutes = int(exam_data.get('duration_minutes') or 0) or exam.duration_minutes or 30
    exam.pass_mark = float(exam_data.get('pass_mark') or 0)
    exam.is_required = bool(exam_data.get('is_required'))
    exam.is_active = bool(exam_data.get('is_active', True))
    exam.allow_retakes = bool(exam_data.get('allow_retakes', True))
    exam.required_to_complete_course = bool(exam_data.get('required_to_complete_course'))

    course_id = exam_data.get('course_id')
    exam.course_id = int(course_id) if course_id else None

    trigger_lesson_id = exam_data.get('trigger_lesson_id')
    exam.trigger_lesson_id = int(trigger_lesson_id) if trigger_lesson_id else None

    settings = exam_data.get('settings') or {}
    exam.settings = settings

    if questions_data is not None:
        existing_questions = {question.id: question for question in exam.questions}
        keep_ids = set()

        for index, question_data in enumerate(questions_data or []):
            question_id = question_data.get('id')
            if question_id and question_id in existing_questions:
                question = existing_questions[question_id]
            else:
                question = ExamQuestion(exam_id=exam.id)
                db.session.add(question)

            question.question_type = question_data.get('question_type') or question_data.get('type') or 'multiple_choice'
            question.text = question_data.get('text', '').strip()
            question.options = question_data.get('options') or []
            question.correct_answers = question_data.get('correct_answers') or []
            question.points = float(question_data.get('points') or 0) or 0.0
            question.is_required = bool(question_data.get('is_required', True))
            question.order_index = index
            question.config = question_data.get('config') or {}

            keep_ids.add(question.id)

        for question in list(existing_questions.values()):
            if question.id not in keep_ids:
                db.session.delete(question)


def _save_exam_payload(exam: Exam | None, exam_data: dict, questions_data: list[dict] | None) -> Exam:
    if exam is None:
        exam = Exam(
            title=exam_data.get('title') or 'Untitled Exam',
            duration_minutes=int(exam_data.get('duration_minutes') or 30)
        )
        db.session.add(exam)
        db.session.flush()

    _apply_exam_updates(exam, exam_data, questions_data)
    db.session.flush()
    return exam


def _next_question_order(exam: Exam) -> int:
    if not exam.questions:
        return 0
    return max((question.order_index or 0) for question in exam.questions) + 1


def _parse_exam_payload():
    payload = request.get_json(silent=True) or {}
    exam_data = payload.get('exam') or {}
    questions_data = payload.get('questions') if 'questions' in payload else None
    return exam_data, questions_data


@app.route('/admin/exams/save', methods=['POST'])
@login_required
@admin_only
def admin_exam_create():
    exam_data, questions_data = _parse_exam_payload()
    exam = _save_exam_payload(None, exam_data, questions_data)
    db.session.commit()
    return jsonify({'success': True, 'exam_id': exam.id})


def _question_payload_from_request() -> dict:
    payload = request.get_json(silent=True) or {}
    question_type = payload.get('question_type') or payload.get('type') or 'multiple_choice'
    text = (payload.get('text') or '').strip()
    raw_points = payload.get('points')
    try:
        points_value = float(raw_points)
    except (TypeError, ValueError):
        points_value = 0.0
    return {
        'question_type': question_type,
        'text': text or 'Untitled question',
        'options': payload.get('options') or [],
        'correct_answers': payload.get('correct_answers') or [],
        'points': points_value,
        'is_required': bool(payload.get('is_required', True)),
        'order_index': payload.get('order_index'),
        'config': payload.get('config') or {}
    }


@app.route('/admin/exams/<int:exam_id>/questions', methods=['POST'])
@login_required
@admin_only
def admin_exam_question_create(exam_id):
    exam = Exam.query.get_or_404(exam_id)
    data = _question_payload_from_request()

    order_index = data['order_index']
    if order_index is None:
        order_index = _next_question_order(exam)

    question = ExamQuestion(
        exam_id=exam.id,
        question_type=data['question_type'],
        text=data['text'],
        options=data['options'],
        correct_answers=data['correct_answers'],
        points=data['points'],
        is_required=data['is_required'],
        order_index=order_index,
        config=data['config']
    )
    db.session.add(question)
    db.session.commit()

    return jsonify({'success': True, 'question': _serialize_exam_question(question)})


@app.route('/admin/exams/<int:exam_id>/questions/<int:question_id>', methods=['PATCH'])
@login_required
@admin_only
def admin_exam_question_update(exam_id, question_id):
    exam = Exam.query.get_or_404(exam_id)
    question = ExamQuestion.query.filter_by(exam_id=exam.id, id=question_id).first_or_404()
    data = _question_payload_from_request()

    question.question_type = data['question_type']
    question.text = data['text']
    question.options = data['options']
    question.correct_answers = data['correct_answers']
    question.points = data['points']
    question.is_required = data['is_required']
    if data['order_index'] is not None:
        question.order_index = int(data['order_index'])
    question.config = data['config']

    db.session.commit()

    return jsonify({'success': True, 'question': _serialize_exam_question(question)})


@app.route('/admin/exams/<int:exam_id>/questions/<int:question_id>', methods=['DELETE'])
@login_required
@admin_only
def admin_exam_question_delete(exam_id, question_id):
    exam = Exam.query.get_or_404(exam_id)
    question = ExamQuestion.query.filter_by(exam_id=exam.id, id=question_id).first_or_404()
    db.session.delete(question)
    db.session.commit()
    return jsonify({'success': True})


@app.route('/admin/exams/<int:exam_id>/questions/reorder', methods=['POST'])
@login_required
@admin_only
def admin_exam_question_reorder(exam_id):
    exam = Exam.query.get_or_404(exam_id)
    payload = request.get_json(silent=True) or {}
    order = payload.get('order') or []

    question_lookup = {question.id: question for question in exam.questions}

    for index, question_id in enumerate(order):
        try:
            resolved_id = int(question_id)
        except (TypeError, ValueError):
            continue
        question = question_lookup.get(resolved_id)
        if question is None:
            continue
        question.order_index = index

    db.session.commit()
    return jsonify({'success': True})


@app.route('/admin/exams/<int:exam_id>/save', methods=['POST'])
@login_required
@admin_only
def admin_exam_update(exam_id):
    exam = Exam.query.get_or_404(exam_id)
    exam_data, questions_data = _parse_exam_payload()
    exam = _save_exam_payload(exam, exam_data, questions_data)
    db.session.commit()
    return jsonify({'success': True, 'exam_id': exam.id})


@app.route('/admin/exams/results/<int:exam_id>')
@login_required
@admin_only
def admin_exam_results(exam_id):
    exam = Exam.query.get_or_404(exam_id)
    attempts = ExamAttempt.query.filter_by(exam_id=exam.id).order_by(ExamAttempt.start_time.desc()).all()
    analytics = _exam_statistics(exam.id)

    question_map = {
        question.id: question
        for question in exam.questions
    }

    if question_map:
        question_stats = db.session.query(
            ExamAnswer.question_id,
            func.count(ExamAnswer.id).label('total'),
            func.sum(case((ExamAnswer.is_correct == True, 1), else_=0)).label('correct')  # noqa: E712
        ).filter(
            ExamAnswer.question_id.in_(question_map.keys())
        ).group_by(ExamAnswer.question_id).all()
    else:
        question_stats = []

    question_analytics = []
    for row in question_stats:
        question = question_map.get(row.question_id)
        total = int(row.total or 0)
        correct = int(row.correct or 0)
        question_analytics.append({
            'question': question,
            'total': total,
            'correct': correct,
            'incorrect': total - correct
        })

    attempts_payload = [_admin_attempt_payload(attempt) for attempt in attempts]

    return render_template(
        'admin_exam_results.html',
        exam=exam,
        course=exam.course,
        analytics=analytics,
        attempts=attempts,
        attempts_payload=attempts_payload,
        question_analytics=question_analytics
    )


@app.route('/admin/exams/attempts/<int:attempt_id>')
@login_required
@admin_only
def admin_exam_attempt_detail(attempt_id):
    attempt = ExamAttempt.query.get_or_404(attempt_id)
    payload = _admin_attempt_payload(attempt)
    return jsonify({'attempt': payload})


@app.route('/admin/exams/attempts/<int:attempt_id>/grade', methods=['POST'])
@login_required
@admin_only
def admin_exam_attempt_grade(attempt_id):
    attempt = ExamAttempt.query.get_or_404(attempt_id)
    payload = request.get_json(silent=True) or {}

    answers_payload = payload.get('answers') or []
    overall_feedback = payload.get('overall_feedback')
    status = payload.get('status')

    answers_lookup = {answer.id: answer for answer in attempt.answers}

    for item in answers_payload:
        answer_id = item.get('id')
        answer = answers_lookup.get(answer_id)
        if not answer:
            continue
        if 'points_awarded' in item:
            try:
                answer.points_awarded = float(item['points_awarded'])
            except (TypeError, ValueError):
                answer.points_awarded = answer.points_awarded or 0.0
        if 'is_correct' in item:
            value = item['is_correct']
            if value is None:
                answer.is_correct = None
            else:
                answer.is_correct = bool(value)
        if 'feedback' in item:
            answer.feedback = item['feedback']

    attempt.overall_feedback = overall_feedback
    target_status = status or 'graded'

    _compute_attempt_score(attempt)
    attempt.status = target_status

    db.session.commit()

    return jsonify({'success': True, 'attempt': _admin_attempt_payload(attempt)})
@app.route('/admin/users')
@login_required
@admin_only
def admin_users():
    users = User.query.order_by(User.id).all()
    courses = Course.query.order_by(Course.year, Course.name).all()
    return render_template("admin_users.html", users=users, courses=courses)


@app.route('/admin/email', methods=['GET', 'POST'])
@login_required
@admin_only
def admin_email():
    """Admin bulk email sending interface"""
    if request.method == 'GET':
        users = User.query.order_by(User.username).all()
        courses = Course.query.order_by(Course.name).all()
        return render_template('admin_email.html', users=users, courses=courses)

    # POST: Send bulk email
    subject = request.form.get('subject', '').strip()
    message_html = request.form.get('message', '').strip()
    recipient_type = request.form.get('recipient_type', 'selected')

    if not subject or not message_html:
        flash('Subject and message are required', 'error')
        return redirect(url_for('admin_email'))

    # Collect recipients based on selection
    recipients = []

    if recipient_type == 'all':
        recipients = User.query.all()
    elif recipient_type == 'course':
        course_name = request.form.get('course_filter', '').strip().lower()
        if course_name:
            recipients = User.query.filter(
                func.lower(User.courses).like(f'%{course_name}%')
            ).all()
    elif recipient_type == 'selected':
        user_ids = request.form.getlist('user_ids[]')
        if user_ids:
            recipients = User.query.filter(User.id.in_(user_ids)).all()

    if not recipients:
        flash('No recipients selected', 'error')
        return redirect(url_for('admin_email'))

    # Wrap message in base email template
    from flask import render_template_string
    html_body = render_template('emails/base_email.html')
    # Replace the content block with custom message
    html_body = html_body.replace('{% block content %}', '').replace('{% endblock %}', message_html)

    # Send emails
    from email_utils import send_bulk_email
    recipient_emails = [user.email for user in recipients if user.email]

    success_count = send_bulk_email(
        recipients=recipient_emails,
        subject=subject,
        html_body=message_html,
        text_body=None
    )

    flash(f'Successfully queued {success_count} emails out of {len(recipient_emails)} recipients', 'success')
    return redirect(url_for('admin_email'))


@app.route('/admin/test-email', methods=['GET', 'POST'])
@login_required
@admin_only
def admin_test_email():
    """Send test email to admin"""
    if request.method == 'GET':
        admin_user = User.query.filter_by(username=session['user']).first()
        return render_template('admin_test_email.html', admin_email=admin_user.email if admin_user else '')

    # POST: Send test email
    test_email = request.form.get('email', '').strip()

    if not test_email:
        flash('Email address is required', 'error')
        return redirect(url_for('admin_test_email'))

    from email_utils import send_email

    html_body = render_template('emails/base_email.html').replace(
        '{% block content %}',
        '<h2>Test Email</h2><p>This is a test email from Al-Baqi Academy email system.</p><p>If you received this, your email configuration is working correctly!</p>'
    ).replace('{% endblock %}', '')

    success = send_email(
        to=test_email,
        subject='Test Email - Al-Baqi Academy',
        html_body=html_body,
        async_send=False
    )

    if success:
        flash(f'Test email sent to {test_email}', 'success')
    else:
        flash('Failed to send test email. Check server logs.', 'error')

    return redirect(url_for('admin_test_email'))


@app.route('/admin/terms', methods=['GET', 'POST'])
@login_required
@admin_only
def admin_terms():
    """Admin interface to edit Terms of Service"""
    if request.method == 'GET':
        terms_setting = SiteSetting.query.filter_by(key='terms_of_service').first()
        terms_content = terms_setting.value if terms_setting else ""
        return render_template('admin_terms.html', terms_content=terms_content)

    # POST: Save updated terms
    new_terms = request.form.get('terms_content', '').strip()

    terms_setting = SiteSetting.query.filter_by(key='terms_of_service').first()
    if terms_setting:
        terms_setting.value = new_terms
        terms_setting.updated_at = datetime.now(timezone.utc)
    else:
        terms_setting = SiteSetting(
            key='terms_of_service',
            value=new_terms,
            updated_at=datetime.now(timezone.utc)
        )
        db.session.add(terms_setting)

    db.session.commit()
    flash('Terms of Service updated successfully!', 'success')
    return redirect(url_for('admin_terms'))


@app.route('/admin/users/<int:user_id>/tracking')
@login_required
@admin_only
def admin_user_tracking(user_id):
    user = User.query.get_or_404(user_id)

    assigned_courses = set(filter(None, (course_name.strip().lower() for course_name in user.get_courses())))
    progress_records = {record.course_id: record for record in user.course_progress}

    tracked_courses = []
    if assigned_courses:
        course_candidates = Course.query.filter(func.lower(Course.name).in_(assigned_courses)).order_by(Course.year, Course.name).all()
    else:
        course_candidates = []

    seen_course_ids = set()
    for course in course_candidates:
        progress_record = progress_records.get(course.id)
        completed_weeks = max(((progress_record.progress if progress_record else 1) - 1), 0)
        total_lessons = len(course.lessons)
        tracked_courses.append({
            "course": course,
            "progress": progress_record,
            "completed_weeks": completed_weeks,
            "total_lessons": total_lessons,
            "next_week": (progress_record.progress if progress_record else 1)
        })
        seen_course_ids.add(course.id)

    for record in progress_records.values():
        if record.course_id in seen_course_ids:
            continue
        course = record.course
        completed_weeks = max((record.progress - 1), 0)
        total_lessons = len(course.lessons) if course else 0
        tracked_courses.append({
            "course": course,
            "progress": record,
            "completed_weeks": completed_weeks,
            "total_lessons": total_lessons,
            "next_week": record.progress
        })

    tracked_courses.sort(key=lambda row: (
        row["course"].year if row["course"] else 0,
        row["course"].name if row["course"] else ""
    ))

    quiz_attempts = QuizAttempt.query.filter_by(user_id=user.id).all()
    quiz_rows = []
    for attempt in quiz_attempts:
        lesson = attempt.lesson or _fetch_lesson(attempt.lesson_id)
        course = lesson.course if lesson else _fetch_course(attempt.course_id)
        try:
            detail_data = json.loads(attempt.detail_json) if attempt.detail_json else []
        except (TypeError, json.JSONDecodeError):
            detail_data = []
        wrong_details = [item for item in detail_data if not item.get('is_correct')]
        total_questions = attempt.total_questions or (len(lesson.quizzes) if lesson else 0)
        score_percent = 0
        if total_questions:
            score_percent = round((attempt.last_score / total_questions) * 100)
        course_label = "Unknown course"
        quiz_label = lesson.title if lesson else f"Quiz #{attempt.id}"
        if course:
            course_label = f"{course.name.title()} (Year {course.year})"
            quiz_label = f"Week {lesson.week}: {lesson.title}" if lesson else quiz_label
        quiz_rows.append({
            "attempt": attempt,
            "lesson": lesson,
            "course": course,
            "score_percent": score_percent,
            "wrong_details": wrong_details,
            "detail_data": detail_data,
            "course_label": course_label,
            "quiz_label": quiz_label,
            "total_questions": total_questions
        })

    quiz_rows.sort(key=lambda row: (
        row["course"].year if row["course"] else 0,
        row["course"].name if row["course"] else "",
        row["lesson"].week if row["lesson"] else 0
    ))

    revision_course_filter = request.args.get('revision_course_id', type=int)
    revision_lesson_filter = request.args.get('revision_lesson_id', type=int)
    revision_sort = request.args.get('revision_sort', 'recent')

    all_retake_attempts = QuizRetakeAttempt.query.filter_by(user_id=user.id).all()

    course_revision_summary = {}
    lesson_options_map = {}
    for attempt in all_retake_attempts:
        course_ref = attempt.course or _fetch_course(attempt.course_id)
        lesson_ref = attempt.lesson or _fetch_lesson(attempt.lesson_id)
        summary_entry = course_revision_summary.setdefault(attempt.course_id, {
            "course": course_ref,
            "count": 0
        })
        summary_entry["count"] += 1
        if course_ref and lesson_ref:
            lesson_options_map.setdefault(course_ref.id, {})[lesson_ref.id] = lesson_ref

    retake_query = QuizRetakeAttempt.query.filter_by(user_id=user.id)
    if revision_course_filter:
        retake_query = retake_query.filter_by(course_id=revision_course_filter)
    if revision_lesson_filter:
        retake_query = retake_query.filter_by(lesson_id=revision_lesson_filter)

    retake_attempts = retake_query.order_by(QuizRetakeAttempt.created_at.desc()).all()

    if revision_sort == 'score_desc':
        retake_attempts.sort(key=lambda a: (a.score / (a.total_questions or 1), a.created_at or datetime.min), reverse=True)
    elif revision_sort == 'score_asc':
        retake_attempts.sort(key=lambda a: (a.score / (a.total_questions or 1), a.created_at or datetime.min))
    elif revision_sort == 'oldest':
        retake_attempts.sort(key=lambda a: a.created_at or datetime.min)

    retake_rows = []
    for attempt in retake_attempts:
        course_ref = attempt.course or _fetch_course(attempt.course_id)
        lesson_ref = attempt.lesson or _fetch_lesson(attempt.lesson_id)
        answers_data = []
        wrong_data = []
        if attempt.is_complete:
            try:
                answers_data = json.loads(attempt.answers_json) if attempt.answers_json else []
            except (TypeError, json.JSONDecodeError):
                answers_data = []
            try:
                wrong_data = json.loads(attempt.wrong_questions_json) if attempt.wrong_questions_json else []
            except (TypeError, json.JSONDecodeError):
                wrong_data = []
        score_percent = 0
        if attempt.total_questions:
            score_percent = round((attempt.score / attempt.total_questions) * 100)
        retake_rows.append({
            "attempt": attempt,
            "course": course_ref,
            "lesson": lesson_ref,
            "answers": answers_data,
            "wrongs": wrong_data,
            "score_percent": score_percent
        })

    student_files = StudentHubFile.query.filter_by(user_id=user.id).order_by(StudentHubFile.uploaded_at.desc()).all()

    activity_log = []
    if user.created_at:
        activity_log.append({
            "timestamp": user.created_at,
            "label": "Signed up",
            "icon": "📝"
        })
    if user.last_login:
        activity_log.append({
            "timestamp": user.last_login,
            "label": "Last login",
            "icon": "🔑"
        })
    for row in quiz_rows:
        attempt = row["attempt"]
        if attempt.last_attempt_at:
            status = "Passed" if attempt.passed else "Attempted"
            activity_log.append({
                "timestamp": attempt.last_attempt_at,
                "label": f"{status} quiz – {row['quiz_label']}",
                "icon": "🧠"
            })
    for upload in student_files:
        activity_log.append({
            "timestamp": upload.uploaded_at,
            "label": f"Uploaded {upload.original_name}",
            "icon": "📁"
        })
    for attempt in all_retake_attempts:
        if attempt.created_at:
            lesson_ref = _fetch_lesson(attempt.lesson_id)
            if lesson_ref:
                lesson_label = f"Week {lesson_ref.week}: {lesson_ref.title}"
            else:
                lesson_label = "Entire course"
            activity_log.append({
                "timestamp": attempt.created_at,
                "label": f"Revision attempt – {lesson_label} ({attempt.score}/{attempt.total_questions})",
                "icon": "📝"
            })

    agreements = CourseAgreement.query.filter_by(user_id=user.id).order_by(CourseAgreement.accepted_at.desc()).all()
    agreement_rows = []
    agreements_by_course_id = {agreement.course_id: agreement for agreement in agreements}

    allowed_course_names = {
        name.strip().lower()
        for name in user.get_courses()
        if name and name.strip()
    }

    course_candidates = []
    if allowed_course_names:
        course_candidates = Course.query.filter(
            func.lower(Course.name).in_(allowed_course_names)
        ).order_by(Course.name.asc(), Course.year.asc()).all()

    seen_course_ids: set[int] = set()
    for course_ref in course_candidates:
        agreement = agreements_by_course_id.get(course_ref.id)
        seen_course_ids.add(course_ref.id)
        agreement_rows.append({
            "course": course_ref,
            "course_label": f"{course_ref.name.title()} (Year {course_ref.year})",
            "accepted_at": agreement.accepted_at if agreement else None,
            "accepted": agreement is not None
        })

    for agreement in agreements:
        if agreement.course_id not in seen_course_ids:
            course_ref = agreement.course or _fetch_course(agreement.course_id)
            label = None
            if course_ref:
                label = f"{course_ref.name.title()} (Year {course_ref.year})"
            else:
                label = f"Course #{agreement.course_id}"
            agreement_rows.append({
                "course": course_ref,
                "course_label": label,
                "accepted_at": agreement.accepted_at,
                "accepted": True
            })

    unmatched_names = allowed_course_names - {
        course.name.strip().lower() for course in course_candidates
    }
    for name_value in sorted(unmatched_names):
        agreement_rows.append({
            "course": None,
            "course_label": name_value.title(),
            "accepted_at": None,
            "accepted": False
        })

    agreement_rows.sort(key=lambda row: (
        (row["course"].name.lower() if row.get("course") else row.get("course_label", "")),
        (row["course"].year if row.get("course") else 0)
    ))

    activity_log = sorted(
        [item for item in activity_log if item.get("timestamp")],
        key=lambda entry: entry["timestamp"],
        reverse=True
    )

    return render_template(
        "admin_user_tracking.html",
        target_user=user,
        tracked_courses=tracked_courses,
        quiz_rows=quiz_rows,
        student_files=student_files,
        activity_log=activity_log,
        retake_rows=retake_rows,
        revision_course_filter=revision_course_filter,
        revision_lesson_filter=revision_lesson_filter,
        revision_sort=revision_sort,
        revision_course_summary=sorted(course_revision_summary.values(), key=lambda entry: (entry["course"].name if entry["course"] else "")),
        revision_lesson_options={cid: sorted(lessons.values(), key=lambda lesson: lesson.week) for cid, lessons in lesson_options_map.items()},
        agreement_rows=agreement_rows
    )


@app.route('/admin/users/<int:user_id>/tracking/reset/<int:lesson_id>', methods=['POST'])
@login_required
@admin_only
def admin_reset_quiz(user_id, lesson_id):
    user = User.query.get_or_404(user_id)
    attempt = QuizAttempt.query.filter_by(user_id=user.id, lesson_id=lesson_id).first()
    redirect_target = request.form.get('redirect')
    if not redirect_target or not redirect_target.startswith('/'):
        redirect_target = url_for('admin_user_tracking', user_id=user.id)

    if not attempt:
        flash("⚠️ No quiz attempts found to reset for this lesson.")
        return redirect(redirect_target)

    attempt.attempt_count = 0
    attempt.last_score = 0
    attempt.best_score = 0
    attempt.correct_count = 0
    attempt.wrong_count = 0
    attempt.passed = False
    attempt.detail_json = None
    attempt.last_attempt_at = None

    lesson = attempt.lesson
    if lesson:
        progress_record = UserCourseProgress.query.filter_by(
            user_id=user.id,
            course_id=lesson.course_id
        ).first()
        if progress_record and progress_record.progress > lesson.week:
            progress_record.progress = max(1, lesson.week)

    db.session.commit()
    flash("🔁 Quiz attempts reset. The student must retake this quiz.")
    return redirect(redirect_target)


@app.route('/admin/users/<int:user_id>/tracking/pass/<int:lesson_id>', methods=['POST'])
@login_required
@admin_only
def admin_force_pass_quiz(user_id, lesson_id):
    user = User.query.get_or_404(user_id)
    lesson = Lesson.query.get_or_404(lesson_id)
    redirect_target = request.form.get('redirect')
    if not redirect_target or not redirect_target.startswith('/'):
        redirect_target = url_for('admin_user_tracking', user_id=user.id)

    attempt = QuizAttempt.query.filter_by(user_id=user.id, lesson_id=lesson.id).first()
    if not attempt:
        attempt = QuizAttempt(user_id=user.id, lesson_id=lesson.id)
        db.session.add(attempt)

    attempt.attempt_count = (attempt.attempt_count or 0) + 1
    total_questions = len(lesson.quizzes)
    attempt.total_questions = total_questions
    attempt.last_score = total_questions
    attempt.best_score = max(attempt.best_score or 0, total_questions)
    attempt.correct_count = total_questions
    attempt.wrong_count = 0
    attempt.passed = True
    attempt.detail_json = None
    attempt.last_attempt_at = datetime.now(timezone.utc)

    progress_record = UserCourseProgress.query.filter_by(
        user_id=user.id,
        course_id=lesson.course_id
    ).first()
    next_week = lesson.week + 1
    if not progress_record:
        progress_record = UserCourseProgress(
            user_id=user.id,
            course_id=lesson.course_id,
            progress=next_week
        )
        db.session.add(progress_record)
    elif progress_record.progress < next_week:
        progress_record.progress = next_week

    db.session.commit()
    flash("✅ Quiz marked as passed for this student.")
    return redirect(redirect_target)


@app.route('/admin/student-hub/delete/<int:file_id>', methods=['POST'])
@login_required
@admin_only
def admin_student_hub_delete(file_id):
    file_record = StudentHubFile.query.get_or_404(file_id)
    _remove_student_hub_file_from_disk(file_record)
    db.session.delete(file_record)
    db.session.commit()
    flash("🗑 Student Hub file deleted.")

    redirect_target = request.form.get('redirect')
    if redirect_target and redirect_target.startswith('/'):
        return redirect(redirect_target)
    return redirect(url_for('admin_student_files'))


@app.route('/admin/student-hub/feedback/<int:file_id>', methods=['POST'])
@login_required
@admin_only
def admin_student_hub_feedback(file_id):
    file_record = StudentHubFile.query.get_or_404(file_id)
    action = request.form.get('action', 'save')

    if action == 'clear':
        file_record.feedback_text = None
        file_record.feedback_grade = None
        db.session.commit()
        flash("📝 Feedback cleared.")
    else:
        feedback_text = request.form.get('feedback_text', '').strip()
        raw_grade = request.form.get('feedback_grade')

        try:
            normalized_grade = _normalize_feedback_grade(raw_grade)
        except ValueError as exc:
            flash(f"⚠️ {exc}")
            redirect_target = request.form.get('redirect')
            if redirect_target and redirect_target.startswith('/'):
                return redirect(redirect_target)
            return redirect(url_for('admin_student_files'))

        file_record.feedback_text = feedback_text or None
        file_record.feedback_grade = normalized_grade
        db.session.commit()
        flash("✅ Feedback saved.")

    redirect_target = request.form.get('redirect')
    if redirect_target and redirect_target.startswith('/'):
        return redirect(redirect_target)
    return redirect(url_for('admin_student_files'))


@app.route('/admin/update_role/<int:user_id>', methods=['POST'])
@login_required
@admin_only
def update_role(user_id):
    user = User.query.get_or_404(user_id)
    new_role = request.form.get("role")

    if new_role in ["user", "paid", "admin"]:
        user.role = new_role
        db.session.commit()
        flash(f"✅ Updated {user.username} to role: {new_role}")
    else:
        flash("⚠️ Invalid role selected")

    redirect_target = request.form.get('redirect')
    if redirect_target and redirect_target.startswith('/'):
        return redirect(redirect_target)
    return redirect(url_for('admin_users'))


@app.route('/admin/update_courses/<int:user_id>', methods=['POST'])
@login_required
@admin_only
def update_courses(user_id):
    user = User.query.get_or_404(user_id)
    selected = request.form.getlist("courses")  # ✅ capture multiple ticked boxes
    user.set_courses(selected)
    db.session.commit()
    flash(f"✅ Updated courses for {user.username}: {', '.join(selected)}")
    redirect_target = request.form.get('redirect')
    if redirect_target and redirect_target.startswith('/'):
        return redirect(redirect_target)
    return redirect(url_for('admin_users'))


@app.route('/admin/add-user', methods=['GET', 'POST'])
@login_required
@admin_only
def admin_add_user():
    if request.method == 'POST':
        username = request.form['username']
        raw_password = request.form['password']
        password = bcrypt.generate_password_hash(raw_password).decode('utf-8')
        full_name = request.form.get('full_name')
        age = request.form.get('age')
        phone_number = request.form.get('phone_number')
        email = request.form.get('email')
        role = request.form.get('role', 'user')  # default is user if not set

        if User.query.filter((User.username == username) | (User.email == email)).first():
            flash("⚠️ Username or Email already exists!")
            return redirect(url_for('admin_add_user'))

        new_user = User(
            username=username,
            password=password,
            full_name=full_name,
            age=age,
            phone_number=phone_number,
            email=email,
            role=role
        )
        db.session.add(new_user)
        db.session.commit()

        flash(f"✅ User {username} added successfully with role {role}")
        return redirect(url_for('admin_add_user'))

    return render_template("admin_add_user.html")

@app.route('/admin/lesson/<int:lesson_id>/delete', methods=['POST'])
@login_required
@admin_only
def delete_lesson(lesson_id):
    lesson = Lesson.query.get_or_404(lesson_id)
    db.session.delete(lesson)
    db.session.commit()
    flash("🗑 Lesson deleted")
    return redirect(url_for('manage_courses'))


@app.route('/admin/delete/<int:user_id>', methods=['POST'])
@login_required
@admin_only
def delete_user(user_id):
    user = User.query.get_or_404(user_id)

    # Prevent deleting yourself
    if user.username == session.get("user"):
        flash("⚠️ You cannot delete your own account while logged in.")
        redirect_target = request.form.get('redirect')
        if redirect_target and redirect_target.startswith('/'):
            return redirect(redirect_target)
        return redirect(url_for('admin_users'))
    
    db.session.delete(user)
    db.session.commit()
    flash(f"Deleted user {user.username}")
    redirect_target = request.form.get('redirect')
    if redirect_target and redirect_target.startswith('/'):
        return redirect(redirect_target)
    return redirect(url_for('admin_users'))

@app.route('/admin/course/<int:course_id>/delete', methods=['POST'])
@login_required
@admin_only
def delete_course(course_id):
    course = Course.query.get_or_404(course_id)
    db.session.delete(course)
    db.session.commit()
    flash(f"🗑 Course {course.name} deleted along with its lessons and quizzes.")
    return redirect(url_for('manage_courses'))


# -------------------- Admin: Testimonials Management -------------------- #
@app.route('/admin/testimonials')
@login_required
@admin_only
def admin_testimonials():
    """Admin page to manage all testimonials."""
    status_filter = request.args.get('status', 'all')

    query = Testimonial.query
    if status_filter and status_filter != 'all':
        query = query.filter_by(status=status_filter)

    testimonials = query.order_by(Testimonial.created_at.desc()).all()

    return render_template(
        'admin_testimonials.html',
        testimonials=testimonials,
        status_filter=status_filter
    )


@app.route('/admin/testimonials/<int:testimonial_id>/approve', methods=['POST'])
@login_required
@admin_only
def approve_testimonial(testimonial_id):
    """Approve a testimonial - will show on homepage."""
    testimonial = Testimonial.query.get_or_404(testimonial_id)
    testimonial.status = 'approved'
    db.session.commit()
    flash("✅ Testimonial approved and will appear on homepage.")
    return redirect(url_for('admin_testimonials'))


@app.route('/admin/testimonials/<int:testimonial_id>/reject', methods=['POST'])
@login_required
@admin_only
def reject_testimonial(testimonial_id):
    """Reject a testimonial - will not show on homepage."""
    testimonial = Testimonial.query.get_or_404(testimonial_id)
    testimonial.status = 'rejected'
    db.session.commit()
    flash("❌ Testimonial rejected.")
    return redirect(url_for('admin_testimonials'))


@app.route('/admin/testimonials/<int:testimonial_id>/edit', methods=['POST'])
@login_required
@admin_only
def edit_testimonial(testimonial_id):
    """Edit testimonial name, rating, or review text."""
    testimonial = Testimonial.query.get_or_404(testimonial_id)

    testimonial.name = request.form.get('name', testimonial.name)
    testimonial.rating = int(request.form.get('rating', testimonial.rating))
    testimonial.review = request.form.get('review', testimonial.review)

    db.session.commit()
    flash("✅ Testimonial updated successfully.")
    return redirect(url_for('admin_testimonials'))


@app.route('/admin/testimonials/<int:testimonial_id>/delete', methods=['POST'])
@login_required
@admin_only
def delete_testimonial(testimonial_id):
    """Delete a testimonial permanently."""
    testimonial = Testimonial.query.get_or_404(testimonial_id)
    db.session.delete(testimonial)
    db.session.commit()
    flash("🗑 Testimonial deleted.")
    return redirect(url_for('admin_testimonials'))


@app.route('/admin/users/<int:user_id>/course/<int:course_id>/mark-finished', methods=['POST'])
@login_required
@admin_only
def mark_course_finished(user_id, course_id):
    """Mark a course as finished for a student, allowing them to leave a review."""
    progress = UserCourseProgress.query.filter_by(
        user_id=user_id,
        course_id=course_id
    ).first()

    if not progress:
        # Create progress record if it doesn't exist
        progress = UserCourseProgress(
            user_id=user_id,
            course_id=course_id,
            progress=1,
            is_finished=True
        )
        db.session.add(progress)
    else:
        progress.is_finished = not progress.is_finished  # Toggle

    db.session.commit()

    status = "finished" if progress.is_finished else "not finished"
    flash(f"✅ Course marked as {status}.")
    return redirect(url_for('admin_user_tracking', user_id=user_id))


# -------------------- Admin Subscription Management Routes -------------------- #
@app.route('/admin/subscriptions')
@login_required
@admin_only
def admin_subscriptions():
    """
    Admin view for managing subscription plans and subscribers.
    """
    plans = SubscriptionPlan.query.all()
    all_courses = Course.query.filter_by(is_published=True).order_by(Course.name).all()

    # Get all active subscriptions with user and plan data
    active_subs = db.session.query(UserSubscription, User, SubscriptionPlan).join(
        User, UserSubscription.user_id == User.id
    ).join(
        SubscriptionPlan, UserSubscription.plan_id == SubscriptionPlan.id
    ).filter(
        UserSubscription.status.in_(['active', 'trialing', 'past_due'])
    ).all()

    return render_template(
        'admin_subscriptions.html',
        plans=plans,
        all_courses=all_courses,
        active_subscriptions=active_subs
    )


@app.route('/admin/subscriptions/plan/create', methods=['POST'])
@login_required
@admin_only
def create_subscription_plan():
    """
    Create a new subscription plan.
    """
    from stripe_helpers import sync_plan_with_stripe

    name = request.form.get('name')
    price = request.form.get('price')
    description = request.form.get('description')
    billing_interval = request.form.get('billing_interval', 'monthly')
    trial_days = request.form.get('trial_days', 7)
    grace_period_days = request.form.get('grace_period_days', 3)
    course_ids = request.form.getlist('course_ids[]')

    if not name or not price:
        flash("⚠️ Name and price are required.")
        return redirect(url_for('admin_subscriptions'))

    try:
        # Convert course_ids to integers
        course_ids = [int(cid) for cid in course_ids if cid]

        plan = SubscriptionPlan(
            name=name,
            price=float(price),
            description=description,
            billing_interval=billing_interval,
            course_ids=course_ids,
            trial_days=int(trial_days),
            grace_period_days=int(grace_period_days),
            is_active=True
        )
        db.session.add(plan)
        db.session.commit()

        # Sync with Stripe
        product_id, price_id = sync_plan_with_stripe(plan)
        plan.stripe_product_id = product_id
        plan.stripe_price_id = price_id
        db.session.commit()

        flash(f"✅ Subscription plan '{name}' created successfully.")

    except Exception as e:
        app.logger.error(f"Failed to create subscription plan: {e}")
        flash(f"⚠️ Error creating plan: {e}")
        db.session.rollback()

    return redirect(url_for('admin_subscriptions'))


@app.route('/admin/subscriptions/plan/<int:plan_id>/edit', methods=['POST'])
@login_required
@admin_only
def edit_subscription_plan(plan_id):
    """
    Edit an existing subscription plan.
    """
    from stripe_helpers import sync_plan_with_stripe

    plan = SubscriptionPlan.query.get_or_404(plan_id)

    plan.name = request.form.get('name', plan.name)
    plan.description = request.form.get('description', plan.description)
    plan.trial_days = int(request.form.get('trial_days', plan.trial_days))
    plan.grace_period_days = int(request.form.get('grace_period_days', plan.grace_period_days))

    course_ids = request.form.getlist('course_ids[]')
    if course_ids:
        plan.course_ids = [int(cid) for cid in course_ids if cid]

    # Check if price changed
    new_price = float(request.form.get('price', plan.price))
    price_changed = new_price != float(plan.price)

    if price_changed:
        plan.price = new_price

    try:
        db.session.commit()

        # Re-sync with Stripe if price changed
        if price_changed:
            product_id, price_id = sync_plan_with_stripe(plan)
            plan.stripe_product_id = product_id
            plan.stripe_price_id = price_id
            db.session.commit()

        flash(f"✅ Plan '{plan.name}' updated successfully.")

    except Exception as e:
        app.logger.error(f"Failed to update subscription plan: {e}")
        flash(f"⚠️ Error updating plan: {e}")
        db.session.rollback()

    return redirect(url_for('admin_subscriptions'))


@app.route('/admin/subscriptions/plan/<int:plan_id>/toggle', methods=['POST'])
@login_required
@admin_only
def toggle_subscription_plan(plan_id):
    """
    Activate or deactivate a subscription plan.
    """
    plan = SubscriptionPlan.query.get_or_404(plan_id)
    plan.is_active = not plan.is_active

    try:
        db.session.commit()
        status = "activated" if plan.is_active else "deactivated"
        flash(f"✅ Plan '{plan.name}' {status}.")
    except Exception as e:
        app.logger.error(f"Failed to toggle plan: {e}")
        flash(f"⚠️ Error: {e}")
        db.session.rollback()

    return redirect(url_for('admin_subscriptions'))


@app.route('/admin/subscriptions/plan/<int:plan_id>/delete', methods=['POST'])
@login_required
@admin_only
def delete_subscription_plan(plan_id):
    """
    Delete a subscription plan (only if no active subscriptions).
    """
    plan = SubscriptionPlan.query.get_or_404(plan_id)

    # Check for active subscriptions
    active_subs = UserSubscription.query.filter_by(plan_id=plan_id).filter(
        UserSubscription.status.in_(['active', 'trialing', 'past_due'])
    ).count()

    if active_subs > 0:
        flash(f"⚠️ Cannot delete plan with {active_subs} active subscriptions. Deactivate it instead.")
        return redirect(url_for('admin_subscriptions'))

    try:
        db.session.delete(plan)
        db.session.commit()
        flash(f"✅ Plan '{plan.name}' deleted.")
    except Exception as e:
        app.logger.error(f"Failed to delete plan: {e}")
        flash(f"⚠️ Error: {e}")
        db.session.rollback()

    return redirect(url_for('admin_subscriptions'))


@app.route('/profile', methods=['GET', 'POST'])
@login_required
def profile():
    user = User.query.filter_by(username=session['user']).first()

    if request.method == 'POST':
        action = request.form.get("action")

        # ---- Update Profile ----
        if action == "update_profile":
            new_full_name = request.form.get("full_name")
            new_age = request.form.get("age")
            new_phone = request.form.get("phone_number")
            new_email = request.form.get("email")

            # check duplicates for email and phone
            if new_email and User.query.filter(User.email == new_email, User.id != user.id).first():
                flash("⚠️ Email already in use by another account.")
                return redirect(url_for('profile'))

            if new_phone and User.query.filter(User.phone_number == new_phone, User.id != user.id).first():
                flash("⚠️ Phone number already in use by another account.")
                return redirect(url_for('profile'))

            user.full_name = new_full_name
            user.age = new_age if new_age else None
            user.phone_number = new_phone
            user.email = new_email

            db.session.commit()
            flash("✅ Profile updated successfully!")

        # ---- Change Password ----
        elif action == "change_password":
            old_password = request.form.get("old_password")
            new_password = request.form.get("new_password")
            confirm_password = request.form.get("confirm_password")

            # verify old password
            if not bcrypt.check_password_hash(user.password, old_password):
                flash("❌ Current password is incorrect.")
                return redirect(url_for('profile'))

            # check new password match
            if new_password != confirm_password:
                flash("❌ New passwords do not match.")
                return redirect(url_for('profile'))

            # update password
            user.password = bcrypt.generate_password_hash(new_password).decode('utf-8')
            db.session.commit()
            flash("✅ Password changed successfully!")

        return redirect(url_for('profile'))

    return render_template("profile.html", user=user)


@app.route('/stats')
@login_required
def stats():
    """User stats dashboard showing lessons and courses completed."""
    user = User.query.filter_by(username=session['user']).first()

    # Count total lessons the user has viewed (based on progress records)
    lessons_completed = UserCourseProgress.query.filter_by(user_id=user.id).count()

    # Get all courses and count how many the user has completed
    # A course is considered "completed" if the user's progress >= total lessons in that course
    all_courses = Course.query.all()
    courses_completed = 0
    total_courses = len(all_courses)

    for course in all_courses:
        total_lessons_in_course = Lesson.query.filter_by(course_id=course.id).count()
        user_progress = UserCourseProgress.query.filter_by(
            user_id=user.id,
            course_id=course.id
        ).first()

        if user_progress and user_progress.progress > total_lessons_in_course:
            courses_completed += 1

    # Calculate total lessons across all courses
    total_lessons = Lesson.query.count()

    # Quiz statistics
    total_quiz_attempts = QuizAttempt.query.filter_by(user_id=user.id).count()
    passed_quizzes = QuizAttempt.query.filter_by(user_id=user.id, passed=True).count()

    # Exam statistics
    exam_attempts = (
        ExamAttempt.query
        .filter_by(user_id=user.id)
        .order_by(ExamAttempt.end_time.desc(), ExamAttempt.start_time.desc())
        .all()
    )
    total_exam_attempts = len(exam_attempts)
    passed_exams = sum(1 for attempt in exam_attempts if attempt.passed is True)

    exam_percentages = []
    exam_details = []
    for attempt in exam_attempts:
        if attempt.max_score:
            percentage = round((attempt.score / attempt.max_score) * 100)
            exam_percentages.append(percentage)
        else:
            percentage = 0

        exam_details.append({
            'exam': attempt.exam,
            'course': attempt.course or (attempt.exam.course if attempt.exam else None),
            'attempt': attempt,
            'percentage': percentage
        })

    avg_exam_percentage = round(sum(exam_percentages) / len(exam_percentages), 1) if exam_percentages else None

    return render_template(
        "user/stats.html",
        user=user,
        lessons_completed=lessons_completed,
        total_lessons=total_lessons,
        courses_completed=courses_completed,
        total_courses=total_courses,
        total_quiz_attempts=total_quiz_attempts,
        passed_quizzes=passed_quizzes,
        total_exam_attempts=total_exam_attempts,
        passed_exams=passed_exams,
        avg_exam_percentage=avg_exam_percentage,
        exam_details=exam_details
    )


# -------------------- Testimonials / Reviews -------------------- #
@app.route('/student/review/<int:course_id>', methods=['GET', 'POST'])
@login_required
def submit_review(course_id):
    """Student review submission form - only accessible if course is marked as finished."""
    user = User.query.filter_by(username=session['user']).first()
    course = Course.query.get_or_404(course_id)

    # Check if user has finished this course
    progress = UserCourseProgress.query.filter_by(
        user_id=user.id,
        course_id=course_id
    ).first()

    if not progress or not progress.is_finished:
        flash("❌ You can only review courses you've completed.")
        return redirect(url_for('courses_dashboard'))

    # Check if user already submitted a review for this course
    existing_review = Testimonial.query.filter_by(
        user_id=user.id,
        course_id=course_id
    ).first()

    if existing_review:
        flash("✅ You've already submitted a review for this course.")
        return redirect(url_for('courses_dashboard'))

    if request.method == 'POST':
        name = request.form.get('name', user.full_name or user.username)
        rating = int(request.form.get('rating', 5))
        review_text = request.form.get('review', '').strip()

        if not review_text:
            flash("❌ Please write a review.")
            return redirect(url_for('submit_review', course_id=course_id))

        if rating < 1 or rating > 5:
            flash("❌ Rating must be between 1 and 5.")
            return redirect(url_for('submit_review', course_id=course_id))

        # Create new testimonial
        testimonial = Testimonial(
            user_id=user.id,
            course_id=course_id,
            name=name,
            rating=rating,
            review=review_text,
            status='pending'
        )

        db.session.add(testimonial)
        db.session.commit()

        flash("✅ Thank you! Your review has been submitted and is pending approval.")
        return redirect(url_for('courses_dashboard'))

    return render_template(
        'student_review_form.html',
        user=user,
        course=course
    )


# -------------------- Q&A Chat (User side) -------------------- #
@app.route("/qa")
@login_required
def qa_dashboard():
    user = User.query.filter_by(username=session['user']).first()
    questions = Question.query.filter_by(user_id=user.id).all()
    return render_template("qa_user.html", questions=questions)

@app.route("/qa/new", methods=["POST"])
@login_required
def new_question():
    title = request.form.get("title")
    is_public = True if request.form.get("is_public") == "on" else False
    user = User.query.filter_by(username=session['user']).first()
    q = Question(user_id=user.id, title=title, is_public=is_public)
    db.session.add(q)
    db.session.commit()
    flash("✅ Question created, you can now chat with admin.")
    return redirect(url_for("qa_dashboard"))

@app.route("/qa/<int:question_id>/send", methods=["POST"])
@login_required
def send_message(question_id):
    body = request.form.get("body")
    user = User.query.filter_by(username=session['user']).first()
    q = Question.query.get_or_404(question_id)
    if q.user_id != user.id:
        flash("⚠️ Not your question.")
        return redirect(url_for("qa_dashboard"))
    msg = Message(question_id=question_id, sender="user", body=body)
    db.session.add(msg)
    db.session.commit()
    return redirect(url_for("qa_dashboard"))

@app.route("/qa/message/<int:message_id>/delete", methods=["POST"])
@login_required
def user_delete_message(message_id):
    msg = Message.query.get_or_404(message_id)
    user = User.query.filter_by(username=session['user']).first()

    # only allow the owner of the question to delete their messages
    if msg.sender != "user" or msg.question.user_id != user.id:
        flash("⚠️ You can only delete your own messages.")
        return redirect(url_for("qa_dashboard"))

    db.session.delete(msg)
    db.session.commit()
    flash("🗑 Message deleted.")
    return redirect(url_for("qa_dashboard"))

@app.route("/qa/<int:question_id>/toggle", methods=["POST"])
@login_required
def user_toggle_question(question_id):
    q = Question.query.get_or_404(question_id)
    user = User.query.filter_by(username=session['user']).first()

    if q.user_id != user.id:
        flash("⚠️ You can only change visibility for your own questions.")
        return redirect(url_for("qa_dashboard"))

    q.is_public = not q.is_public
    db.session.commit()
    flash(f"✅ Question changed to {'Public' if q.is_public else 'Private'}.")
    return redirect(url_for("qa_dashboard"))

@app.route("/qa/<int:question_id>/delete", methods=["POST"])
@login_required
def user_delete_question(question_id):
    q = Question.query.get_or_404(question_id)
    user = User.query.filter_by(username=session['user']).first()

    if q.user_id != user.id:
        flash("⚠️ You can only delete your own questions.")
        return redirect(url_for("qa_dashboard"))

    db.session.delete(q)
    db.session.commit()
    flash("🗑 Question deleted.")
    return redirect(url_for("qa_dashboard"))


# -------------------- Q&A Chat (Admin side) -------------------- #
@app.route("/admin/qna")
@login_required
@admin_only
def admin_qna():
    questions = Question.query.all()
    return render_template("admin_qna.html", questions=questions)

@app.route("/admin/qna/<int:question_id>/reply", methods=["POST"])
@login_required
@admin_only
def admin_reply(question_id):
    body = request.form.get("body")
    q = Question.query.get_or_404(question_id)
    msg = Message(question_id=question_id, sender="admin", body=body)
    db.session.add(msg)
    db.session.commit()
    flash("💬 Reply sent.")
    return redirect(url_for("admin_qna"))
# -------------------- Admin Q&A Extras -------------------- #
@app.route("/admin/qna/<int:question_id>/toggle", methods=["POST"])
@login_required
@admin_only
def toggle_question(question_id):
    q = Question.query.get_or_404(question_id)
    q.is_public = not q.is_public
    db.session.commit()
    flash(f"✅ Question visibility changed to {'Public' if q.is_public else 'Private'}")
    return redirect(url_for("admin_qna"))


# -------------------- Admin Q&A Delete -------------------- #
@app.route("/admin/qna/<int:question_id>/delete", methods=["POST"])
@login_required
@admin_only
def delete_question(question_id):
    q = Question.query.get_or_404(question_id)
    db.session.delete(q)
    db.session.commit()
    flash("🗑 Question deleted.")
    return redirect(url_for("admin_qna"))

@app.route("/admin/qna/message/<int:message_id>/delete", methods=["POST"])
@login_required
@admin_only
def delete_message(message_id):
    msg = Message.query.get_or_404(message_id)
    q_id = msg.question_id
    db.session.delete(msg)
    db.session.commit()
    flash("🗑 Message deleted.")
    return redirect(url_for("admin_qna"))

# -------------------- Public Q&A (Homepage embed) -------------------- #
@app.route("/qa/public")
def public_questions():
    questions = Question.query.filter_by(is_public=True).all()
    return render_template("qa_public.html", questions=questions)

@app.route('/admin/course/add', methods=['POST'])
@login_required
@admin_only
def add_course():
    name = request.form.get("name")
    year = int(request.form.get("year", 1))
    description = request.form.get("description")

    # ✅ NEW: Hierarchy and monetization fields
    course_type = request.form.get("course_type", "standalone")
    parent_id = request.form.get("parent_id")
    order_index = int(request.form.get("order_index", 0))
    # Default to True if checkbox not present (for backward compatibility)
    is_published = request.form.get("is_published", "on") == "on"
    show_on_homepage = request.form.get("show_on_homepage") == "on"
    course_assignment = request.form.get("course_assignment", "standalone")
    price = request.form.get("price")

    if not name:
        flash("⚠️ Please provide a course name.")
        return redirect(url_for("manage_courses"))

    # Convert price to decimal or None
    if price and price.strip():
        try:
            price = float(price)
        except ValueError:
            price = None
    else:
        price = None

    # Validate parent_id if sub_course
    if course_type == "sub_course" and parent_id:
        parent_id = int(parent_id)
    else:
        parent_id = None

    course = Course(
        name=name,
        year=year,
        description=description,
        course_type=course_type,
        parent_id=parent_id,
        order_index=order_index,
        is_published=is_published,
        show_on_homepage=show_on_homepage,
        course_assignment=course_assignment,
        price=price
    )
    db.session.add(course)
    db.session.commit()

    flash(f"✅ Course '{name}' added successfully.")
    return redirect(url_for("manage_courses"))

@app.route('/admin/course/<int:course_id>/edit', methods=['POST'])
@login_required
@admin_only
def edit_course(course_id):
    course = Course.query.get_or_404(course_id)
    course.name = request.form.get("name")
    course.year = int(request.form.get("year", 1))
    course.description = request.form.get("description")

    # ✅ NEW: Update hierarchy and monetization fields
    course.course_type = request.form.get("course_type", course.course_type)
    parent_id = request.form.get("parent_id")
    course.order_index = int(request.form.get("order_index", 0))
    course.is_published = request.form.get("is_published") == "on"
    course.show_on_homepage = request.form.get("show_on_homepage") == "on"
    course.course_assignment = request.form.get("course_assignment", course.course_assignment or "standalone")

    # Handle price
    price = request.form.get("price")
    if price and price.strip():
        try:
            course.price = float(price)
        except ValueError:
            course.price = None
    else:
        course.price = None

    # Validate parent_id if sub_course
    if course.course_type == "sub_course" and parent_id:
        course.parent_id = int(parent_id)
    else:
        course.parent_id = None

    db.session.commit()
    flash("✅ Course updated successfully")
    return redirect(url_for('manage_courses'))


@app.route('/admin/course/<int:course_id>/sync-stripe', methods=['POST'])
@login_required
@admin_only
def sync_course_stripe(course_id):
    """Sync a course with Stripe (create/update product and price)"""
    from stripe_helpers import sync_course_with_stripe

    course = Course.query.get_or_404(course_id)

    if not course.price or course.price <= 0:
        flash("⚠️ Cannot sync free courses with Stripe. Set a price first.")
        return redirect(url_for('manage_courses'))

    try:
        product_id, price_id = sync_course_with_stripe(course)
        course.stripe_product_id = product_id
        course.stripe_price_id = price_id
        db.session.commit()
        flash(f"✅ Course '{course.name}' synced with Stripe successfully!")
    except Exception as e:
        app.logger.error(f"Failed to sync course {course_id} with Stripe: {e}")
        flash(f"⚠️ Failed to sync with Stripe: {str(e)}")

    return redirect(url_for('manage_courses'))


# -------------------- Course Access Decorator -------------------- #
def course_required(course_name):
    def decorator(f):
        @wraps(f)
        def decorated_function(*args, **kwargs):
            user = User.query.filter_by(username=session['user']).first()

            # Check access via new hierarchy system or legacy system
            course = Course.query.filter_by(name=course_name).first()
            if course:
                has_access = course.user_has_access(user.id) or course_name in user.get_courses()
            else:
                has_access = course_name in user.get_courses()

            if not has_access and user.role != 'admin':
                flash("⚠️ You don't have access to this course.")
                return redirect(url_for('courses_dashboard'))
            return f(*args, **kwargs)
        return decorated_function
    return decorator


# -------------------- Quiz Submission -------------------- #
@app.route('/submit_quiz/<int:lesson_id>', methods=['POST'])
@login_required
def submit_quiz(lesson_id):
    user = User.query.filter_by(username=session['user']).first()
    lesson = Lesson.query.get_or_404(lesson_id)
    course = lesson.course  # get related course

    # 🔒 Ensure user has access to this course (check both new hierarchy system and legacy)
    has_access = course.user_has_access(user.id) or course.name in user.get_courses()
    if not has_access and user.role != 'admin':
        flash("⚠️ You don't have access to this course.")
        return redirect(url_for('courses_dashboard'))

    # ✅ Mark answers
    correct = 0
    total = len(lesson.quizzes)
    details = []
    for quiz in lesson.quizzes:
        answer = request.form.get(f"q{quiz.id}")
        if answer == quiz.correct_answer:
            correct += 1
        details.append({
            "question": quiz.question,
            "options": {
                "A": quiz.option_a,
                "B": quiz.option_b,
                "C": quiz.option_c
            },
            "correct": quiz.correct_answer,
            "selected": answer,
            "is_correct": answer == quiz.correct_answer
        })

    required_to_pass = 0
    if total > 0:
        required_to_pass = max(1, min(total, round((total * 2) / 3)))

    passed = total == 0 or correct >= required_to_pass

    # ✅ Update dynamic progress only if passed
    if passed and total > 0:
        progress_record = UserCourseProgress.query.filter_by(
            user_id=user.id,
            course_id=course.id
        ).first()

        if not progress_record:
            progress_record = UserCourseProgress(
                user_id=user.id,
                course_id=course.id,
                progress=1
            )
            db.session.add(progress_record)

        if progress_record.progress < lesson.week + 1:
            progress_record.progress = lesson.week + 1

        flash(f"✅ Quiz passed! ({correct}/{total}) Progress unlocked (Week {lesson.week + 1})")
    else:
        flash(f"❌ Quiz requires {required_to_pass}/{total} correct. You scored {correct}. Try again.")

    feedback_payload = {
        "lesson_id": lesson.id,
        "course_id": course.id,
        "score": correct,
        "total": total,
        "required": required_to_pass,
        "passed": passed,
        "details": details
    }
    session['quiz_feedback'] = feedback_payload

    if passed:
        history = session.get('quiz_history', {})
        history[str(lesson.id)] = feedback_payload
        session['quiz_history'] = history
    else:
        history = session.get('quiz_history', {})
        if str(lesson.id) in history:
            history.pop(str(lesson.id), None)
            session['quiz_history'] = history

    attempt_record = QuizAttempt.query.filter_by(user_id=user.id, lesson_id=lesson.id).first()
    if not attempt_record:
        attempt_record = QuizAttempt(user_id=user.id, lesson_id=lesson.id)
        db.session.add(attempt_record)

    attempt_record.attempt_count = (attempt_record.attempt_count or 0) + 1
    attempt_record.last_score = correct
    attempt_record.total_questions = total
    attempt_record.correct_count = correct
    attempt_record.wrong_count = max(total - correct, 0)
    attempt_record.passed = passed
    attempt_record.detail_json = json.dumps(details)
    attempt_record.last_attempt_at = datetime.now(timezone.utc)
    if (attempt_record.best_score or 0) < correct:
        attempt_record.best_score = correct

    db.session.commit()

    if passed:
        return redirect(url_for('course_page', course_name=course.name, year=course.year))

    # On failure, keep the lesson open for retry
    return redirect(url_for(
        'course_page',
        course_name=course.name,
        year=course.year,
        open=lesson.id
    ))



@app.route('/courses/<int:course_id>/quiz-retake/save', methods=['POST'])
@login_required
def course_quiz_retake_save(course_id):
    user = User.query.filter_by(username=session['user']).first()
    course = Course.query.get_or_404(course_id)

    payload = request.get_json(silent=True) or {}
    attempt_id = payload.get('attempt_id')
    quiz_id = payload.get('quiz_id')
    answer = payload.get('answer')

    if not attempt_id or quiz_id is None:
        return {"status": "error", "message": "Missing data"}, 400

    try:
        quiz_id = int(quiz_id)
    except (TypeError, ValueError):
        return {"status": "error", "message": "Invalid quiz id"}, 400

    attempt = QuizRetakeAttempt.query.filter_by(
        id=attempt_id,
        user_id=user.id,
        course_id=course.id,
        is_complete=False
    ).first()
    if not attempt:
        return {"status": "error", "message": "Attempt not found"}, 404

    try:
        question_order = json.loads(attempt.question_order_json or '[]')
    except (TypeError, json.JSONDecodeError):
        question_order = []

    quiz_ids = {item.get('quiz_id') if isinstance(item, dict) else item for item in question_order}
    if quiz_id not in quiz_ids:
        return {"status": "error", "message": "Quiz not part of this attempt"}, 400

    try:
        answers_map = json.loads(attempt.answers_json or '{}')
        if not isinstance(answers_map, dict):
            answers_map = {}
    except (TypeError, json.JSONDecodeError):
        answers_map = {}

    answers_map[str(quiz_id)] = answer
    attempt.answers_json = json.dumps(answers_map)
    attempt.current_index = max(attempt.current_index or 0, len({k for k, v in answers_map.items() if v}))
    if not attempt.created_at:
        attempt.created_at = datetime.now(timezone.utc)

    db.session.commit()
    return {"status": "ok"}


@app.route('/course/<string:course_name>/<int:year>/agreement', methods=['POST'])
@login_required
def accept_course_agreement(course_name, year):
    user = User.query.filter_by(username=session['user']).first()
    course = Course.query.filter_by(name=course_name, year=year).first_or_404()

    # Check access via new hierarchy system or legacy system
    has_access = course.user_has_access(user.id) or course.name in user.get_courses()
    if user.role != 'admin' and not has_access:
        return jsonify({"accepted": False, "message": "Access denied."}), 403

    existing = CourseAgreement.query.filter_by(user_id=user.id, course_id=course.id).first()
    timestamp = datetime.now(timezone.utc)
    if existing:
        existing.accepted_at = timestamp
    else:
        agreement = CourseAgreement(user_id=user.id, course_id=course.id, accepted_at=timestamp)
        db.session.add(agreement)

    db.session.commit()
    return jsonify({"accepted": True, "accepted_at": timestamp.isoformat()})


@app.route('/courses/<int:course_id>/quiz-retake', methods=['GET', 'POST'])
@login_required
def course_quiz_retake(course_id):
    user = User.query.filter_by(username=session['user']).first()
    course = Course.query.get_or_404(course_id)

    # Check access via new hierarchy system or legacy system
    has_access = course.user_has_access(user.id) or course.name in user.get_courses()

    if not has_access:
        # Fallback: check if user has progress record (legacy support)
        progress_record = UserCourseProgress.query.filter_by(
            user_id=user.id,
            course_id=course.id
        ).first()
        has_access = progress_record is not None

    if not has_access and user.role != 'admin':
        flash("⚠️ You don’t have access to this course.")
        return redirect(url_for('courses_dashboard'))

    lessons = Lesson.query.filter_by(course_id=course.id).order_by(Lesson.week).all()
    lessons_with_quiz = [lesson for lesson in lessons if lesson.quizzes]

    if not lessons_with_quiz:
        flash("⚠️ There are no quizzes available for this course yet.")
        return redirect(url_for('course_page', course_name=course.name, year=course.year))

    lessons_lookup = {lesson.id: lesson for lesson in lessons_with_quiz}

    discard_attempt_id = request.args.get('discard_attempt', type=int)
    if discard_attempt_id:
        attempt = QuizRetakeAttempt.query.filter_by(
            id=discard_attempt_id,
            user_id=user.id,
            course_id=course.id,
            is_complete=False
        ).first()
        if attempt:
            db.session.delete(attempt)
            db.session.commit()
            flash("🗑 Unfinished revision attempt discarded.")
        return redirect(url_for('course_quiz_retake', course_id=course.id))

    scope_param = None
    mode_param = None
    selected_lesson_id = None
    result_payload = None
    autosave_attempt = None
    existing_answers = {}
    quiz_items = []
    question_ids_for_form = []
    resume_prompt_attempt = None

    if request.method == 'POST':
        scope_param = request.form.get('scope', 'lesson')
        mode_param = request.form.get('mode', 'list')
    else:
        scope_param = request.args.get('scope') or 'lesson'
        mode_param = request.args.get('mode', 'list')

    if request.method == 'POST' and scope_param == 'lesson':
        selected_lesson_id = request.form.get('lesson_id', type=int)
        selected_lesson = lessons_lookup.get(selected_lesson_id)
        if not selected_lesson:
            flash("⚠️ Invalid lesson selected for quiz retake.")
            return redirect(url_for('course_quiz_retake', course_id=course.id))

        question_id_blob = request.form.get('question_ids', '')
        try:
            question_ids = [int(value) for value in question_id_blob.split(',') if value]
        except ValueError:
            question_ids = []

        quiz_map = {quiz.id: quiz for quiz in selected_lesson.quizzes}

        total_questions = len(question_ids)
        correct_count = 0
        answers_log = []
        wrong_log = []

        for quiz_id in question_ids:
            quiz = quiz_map.get(quiz_id)
            if not quiz:
                continue
            selected_answer = request.form.get(f'quiz_{quiz_id}')
            is_correct = selected_answer == quiz.correct_answer
            if is_correct:
                correct_count += 1
            option_map = {
                'A': quiz.option_a,
                'B': quiz.option_b,
                'C': quiz.option_c
            }
            answer_payload = {
                "quiz_id": quiz_id,
                "question": quiz.question,
                "selected": selected_answer,
                "correct": quiz.correct_answer,
                "options": option_map,
                "is_correct": is_correct
            }
            answers_log.append(answer_payload)
            if not is_correct:
                wrong_log.append(answer_payload)

        attempt_number = QuizRetakeAttempt.query.filter_by(
            user_id=user.id,
            course_id=course.id,
            lesson_id=selected_lesson.id
        ).count() + 1

        attempt_type = 'lesson_random' if mode_param == 'random' else 'lesson_list'

        retake_attempt = QuizRetakeAttempt(
            user_id=user.id,
            course_id=course.id,
            lesson_id=selected_lesson.id,
            attempt_number=attempt_number,
            score=correct_count,
            total_questions=total_questions,
            is_randomized=(mode_param == 'random'),
            attempt_type=attempt_type,
            is_complete=True,
            answers_json=json.dumps(answers_log),
            wrong_questions_json=json.dumps(wrong_log)
        )
        db.session.add(retake_attempt)
        db.session.commit()

        result_payload = {
            "lesson": selected_lesson,
            "total": total_questions,
            "score": correct_count,
            "answers": answers_log,
            "wrongs": wrong_log,
            "attempt": retake_attempt
        }

    elif request.method == 'POST' and scope_param == 'course':
        attempt_id = request.form.get('attempt_id', type=int)
        attempt = QuizRetakeAttempt.query.filter_by(
            id=attempt_id,
            user_id=user.id,
            course_id=course.id,
            is_complete=False
        ).first()
        if not attempt:
            flash("⚠️ Unable to resume that revision attempt.")
            return redirect(url_for('course_quiz_retake', course_id=course.id))

        try:
            question_order = json.loads(attempt.question_order_json or '[]')
        except (TypeError, json.JSONDecodeError):
            question_order = []

        answers_map = {}
        try:
            raw_answers = json.loads(attempt.answers_json or '{}')
            if isinstance(raw_answers, dict):
                answers_map = {str(k): v for k, v in raw_answers.items()}
            elif isinstance(raw_answers, list):
                for item in raw_answers:
                    if isinstance(item, dict):
                        answers_map[str(item.get('quiz_id'))] = item.get('selected')
        except (TypeError, json.JSONDecodeError):
            answers_map = {}

        for key in request.form:
            if key.startswith('quiz_'):
                quiz_id = key.split('_', 1)[1]
                answers_map[quiz_id] = request.form.get(key)

        answers_log = []
        wrong_log = []
        correct_count = 0
        total_questions = 0

        for item in question_order:
            if isinstance(item, dict):
                quiz_id = item.get('quiz_id')
                lesson_id = item.get('lesson_id')
            else:
                quiz_id = item
                lesson_id = None
            quiz = _fetch_quiz(quiz_id)
            if not quiz:
                continue
            total_questions += 1
            selected_answer = answers_map.get(str(quiz_id))
            is_correct = selected_answer == quiz.correct_answer if selected_answer else False
            if is_correct:
                correct_count += 1
            option_map = {
                'A': quiz.option_a,
                'B': quiz.option_b,
                'C': quiz.option_c
            }
            detail_payload = {
                "quiz_id": quiz_id,
                "lesson_id": lesson_id,
                "question": quiz.question,
                "selected": selected_answer,
                "correct": quiz.correct_answer,
                "options": option_map,
                "is_correct": is_correct
            }
            answers_log.append(detail_payload)
            if not is_correct:
                wrong_log.append(detail_payload)

        attempt.score = correct_count
        attempt.total_questions = total_questions
        attempt.answers_json = json.dumps(answers_log)
        attempt.wrong_questions_json = json.dumps(wrong_log)
        attempt.is_complete = True
        attempt.current_index = total_questions
        attempt.is_randomized = attempt.attempt_type.endswith('random')
        attempt.created_at = datetime.now(timezone.utc)
        db.session.commit()

        result_payload = {
            "lesson": None,
            "total": total_questions,
            "score": correct_count,
            "answers": answers_log,
            "wrongs": wrong_log,
            "attempt": attempt,
            "scope": "course"
        }

    pending_attempt = QuizRetakeAttempt.query.filter_by(
        user_id=user.id,
        course_id=course.id,
        is_complete=False
    ).first()

    resume_attempt_id = request.args.get('resume_attempt', type=int)
    active_scope = scope_param
    selected_lesson = None

    if active_scope == 'lesson':
        selected_lesson_id = request.args.get('lesson_id', type=int) if request.method == 'GET' else selected_lesson_id
        selected_lesson = lessons_lookup.get(selected_lesson_id) if selected_lesson_id else None
        if selected_lesson:
            if mode_param == 'random':
                quiz_order = list(selected_lesson.quizzes)
                random.shuffle(quiz_order)
            else:
                quiz_order = list(selected_lesson.quizzes)
                quiz_order.sort(key=lambda q: q.id)

            for quiz in quiz_order:
                quiz_items.append({
                    "id": quiz.id,
                    "question": quiz.question,
                    "options": {
                        'A': quiz.option_a,
                        'B': quiz.option_b,
                        'C': quiz.option_c
                    }
                })
                question_ids_for_form.append(str(quiz.id))

    else:
        active_scope = 'course'
        if resume_attempt_id:
            autosave_attempt = QuizRetakeAttempt.query.filter_by(
                id=resume_attempt_id,
                user_id=user.id,
                course_id=course.id,
                is_complete=False
            ).first()
            if not autosave_attempt:
                flash("⚠️ Could not resume that revision attempt.")
                return redirect(url_for('course_quiz_retake', course_id=course.id))
        elif pending_attempt and not result_payload:
            resume_prompt_attempt = pending_attempt

        if not autosave_attempt and not resume_prompt_attempt and not result_payload:
            attempt_type = 'course_random' if mode_param == 'random' else 'course_list'
            question_order = []
            for lesson in lessons_with_quiz:
                for quiz in lesson.quizzes:
                    question_order.append({
                        "quiz_id": quiz.id,
                        "lesson_id": lesson.id
                    })
            if mode_param == 'random':
                random.shuffle(question_order)

            autosave_attempt = QuizRetakeAttempt(
                user_id=user.id,
                course_id=course.id,
                lesson_id=None,
                attempt_number=QuizRetakeAttempt.query.filter_by(user_id=user.id, course_id=course.id, attempt_type=attempt_type).count() + 1,
                attempt_type=attempt_type,
                is_complete=False,
                is_randomized=(mode_param == 'random')
            )
            autosave_attempt.question_order_json = json.dumps(question_order)
            autosave_attempt.answers_json = json.dumps({})
            db.session.add(autosave_attempt)
            db.session.commit()
            pending_attempt = autosave_attempt

        if autosave_attempt:
            try:
                question_order = json.loads(autosave_attempt.question_order_json or '[]')
            except (TypeError, json.JSONDecodeError):
                question_order = []
            try:
                answers_map = json.loads(autosave_attempt.answers_json or '{}')
                if isinstance(answers_map, dict):
                    existing_answers = {int(k): v for k, v in answers_map.items() if k is not None}
                else:
                    existing_answers = {}
            except (TypeError, json.JSONDecodeError):
                existing_answers = {}

            for item in question_order:
                if isinstance(item, dict):
                    quiz_id = item.get('quiz_id')
                    lesson_id = item.get('lesson_id')
                else:
                    quiz_id = item
                    lesson_id = None
                quiz = _fetch_quiz(quiz_id)
                lesson_ref = lessons_lookup.get(lesson_id) if lesson_id else None
                if not quiz:
                    continue
                quiz_items.append({
                    "id": quiz.id,
                    "question": quiz.question,
                    "options": {
                        'A': quiz.option_a,
                        'B': quiz.option_b,
                        'C': quiz.option_c
                    },
                    "lesson": lesson_ref
                })
                question_ids_for_form.append(str(quiz.id))

    history_records = QuizRetakeAttempt.query.filter_by(
        user_id=user.id,
        course_id=course.id
    ).order_by(QuizRetakeAttempt.created_at.desc()).all()

    history_payload = []
    for attempt in history_records:
        lesson_ref = lessons_lookup.get(attempt.lesson_id) or _fetch_lesson(attempt.lesson_id)
        if attempt.is_complete:
            try:
                answers_data = json.loads(attempt.answers_json) if attempt.answers_json else []
            except (TypeError, json.JSONDecodeError):
                answers_data = []
        else:
            answers_data = []
        history_payload.append({
            "attempt": attempt,
            "lesson": lesson_ref,
            "answers": answers_data
        })

    return render_template(
        "quiz_retake.html",
        course=course,
        lessons=lessons_with_quiz,
        selected_lesson=selected_lesson,
        quiz_items=quiz_items,
        view_mode=mode_param,
        question_ids=','.join(question_ids_for_form),
        result_payload=result_payload,
        history_rows=history_payload,
        user=user,
        scope=active_scope,
        autosave_attempt=autosave_attempt,
        existing_answers=existing_answers,
        resume_prompt_attempt=resume_prompt_attempt
    )

@app.route('/courses')
@login_required
def courses_dashboard():
    user = User.query.filter_by(username=session['user']).first()
    allowed_courses = user.get_courses()

    # ✅ Get published courses grouped by assignment
    # Helper function to calculate course progress
    def get_course_data(course):
        has_access = course.user_has_access(user.id) or course.name in allowed_courses
        is_free = course.is_free()
        requires_approval = False  # Can be extended later for approval-required courses

        # Calculate progress if user has access
        progress_percent = 0
        is_started = False
        progress_record = None
        if has_access or is_free:
            progress_record = UserCourseProgress.query.filter_by(
                user_id=user.id,
                course_id=course.id
            ).first()

            if progress_record:
                is_started = progress_record.progress > 0
                # Calculate percentage: (lessons completed / total lessons) × 100
                total_lessons = Lesson.query.filter_by(course_id=course.id).count()
                if total_lessons > 0:
                    progress_percent = int((progress_record.progress / total_lessons) * 100)

        return {
            'course': course,
            'has_access': has_access,
            'is_free': is_free,
            'requires_approval': requires_approval,
            'is_started': is_started,
            'progress_percent': progress_percent,
            'is_child': course.parent_id is not None,
            'children': []
        }

    def build_course_payload(course):
        payload = get_course_data(course)
        if course.children:
            child_payloads = []
            for child in sorted(
                (c for c in course.children if c.is_published),
                key=lambda c: (c.order_index, c.name.lower())
            ):
                child_payloads.append(get_course_data(child))
            payload['children'] = child_payloads
        return payload

    standalone_courses = [
        build_course_payload(course)
        for course in Course.query.filter_by(
            is_published=True,
            parent_id=None,
            course_assignment='standalone'
        ).order_by(Course.order_index, Course.name).all()
    ]

    year_1_courses = [
        build_course_payload(course)
        for course in Course.query.filter_by(
            is_published=True,
            parent_id=None,
            course_assignment='year_1'
        ).order_by(Course.order_index, Course.name).all()
    ]

    year_2_courses = [
        build_course_payload(course)
        for course in Course.query.filter_by(
            is_published=True,
            parent_id=None,
            course_assignment='year_2'
        ).order_by(Course.order_index, Course.name).all()
    ]

    # Get finished courses for review button
    finished_course_ids = set()
    for progress in user.course_progress:
        if progress.is_finished:
            finished_course_ids.add(progress.course_id)

    # Get courses user has already reviewed
    reviewed_course_ids = set(t.course_id for t in user.testimonials)

    return render_template(
        "courses_dashboard.html",
        user=user,
        standalone_courses=standalone_courses,
        year_1_courses=year_1_courses,
        year_2_courses=year_2_courses,
        finished_course_ids=finished_course_ids,
        reviewed_course_ids=reviewed_course_ids
    )

@app.route('/student-hub')
@login_required
def student_hub():
    user = User.query.filter_by(username=session['user']).first()

    if user.role not in ('paid', 'admin'):
        flash("⚠️ Student Hub is available to paid students only.")
        return redirect(url_for('courses_dashboard'))

    student_files = StudentHubFile.query.filter_by(user_id=user.id).order_by(StudentHubFile.uploaded_at.desc()).all()

    return render_template('student_hub.html', student_files=student_files)


@app.route('/student-hub/upload', methods=['POST'])
@login_required
def student_hub_upload():
    user = User.query.filter_by(username=session['user']).first()

    if user.role not in ('paid', 'admin'):
        flash("⚠️ Student Hub is available to paid students only.")
        return redirect(url_for('courses_dashboard'))

    uploaded_file = request.files.get('student_file')
    if not uploaded_file or uploaded_file.filename == '':
        flash('⚠️ No file selected')
        return redirect(url_for('student_hub'))

    filename = secure_filename(uploaded_file.filename)
    user_folder = os.path.join(app.config['STUDENT_HUB_UPLOAD_FOLDER'], str(user.id))
    os.makedirs(user_folder, exist_ok=True)
    file_path = os.path.join(user_folder, filename)
    uploaded_file.save(file_path)

    file_record = StudentHubFile(
        user_id=user.id,
        stored_name=filename,
        original_name=filename,
        file_path=file_path,
        file_size=os.path.getsize(file_path)
    )
    db.session.add(file_record)
    db.session.commit()

    flash('✅ File uploaded successfully!')
    return redirect(url_for('student_hub'))


@app.route('/student-hub/delete/<int:file_id>', methods=['POST'])
@login_required
def student_hub_delete(file_id):
    user = User.query.filter_by(username=session['user']).first()

    if user.role not in ('paid', 'admin'):
        flash("⚠️ Student Hub is available to paid students only.")
        return redirect(url_for('courses_dashboard'))

    file_record = StudentHubFile.query.get_or_404(file_id)

    if user.role != 'admin' and file_record.user_id != user.id:
        flash("⚠️ You can only delete your own uploads.")
        return redirect(url_for('student_hub'))

    _remove_student_hub_file_from_disk(file_record)
    db.session.delete(file_record)
    db.session.commit()
    flash('🗑 File deleted successfully!')
    return redirect(url_for('student_hub'))


@app.route('/student-hub/download/<int:file_id>')
@login_required
def student_hub_download(file_id):
    user = User.query.filter_by(username=session['user']).first()
    file_record = StudentHubFile.query.get_or_404(file_id)

    if user.role != 'admin' and file_record.user_id != user.id:
        flash("⚠️ You can only access your own Student Hub files.")
        return redirect(url_for('student_hub'))

    folder = os.path.join(app.config['STUDENT_HUB_UPLOAD_FOLDER'], str(file_record.user_id))
    target_path = os.path.join(folder, file_record.stored_name)
    if not os.path.exists(target_path):
        flash("⚠️ File is missing from storage.")
        return redirect(url_for('student_hub'))

    return send_from_directory(
        folder,
        file_record.stored_name,
        as_attachment=True,
        download_name=file_record.original_name
    )


def _resolve_exam_context(course_id: int, exam_id: int) -> tuple['Exam', Course | None]:
    exam = Exam.query.get_or_404(exam_id)
    course = _fetch_course(course_id)

    if exam.course_id and course_id != exam.course_id:
        abort(404)

    if not course and exam.course_id:
        course = _fetch_course(exam.course_id)

    return exam, course


def _assert_exam_permissions(user: User, course: Course | None, exam: 'Exam') -> None:
    if user.role == 'admin':
        return

    if not exam.is_active:
        abort(404)

    if exam.course_id:
        if not course:
            abort(404)
        # Check access via new hierarchy system or legacy system
        has_access = course.user_has_access(user.id) or course.name in user.get_courses()
        if not has_access:
            abort(403)
    else:
        # Standalone exams default to paid or admin users
        if user.role not in {'paid', 'admin'}:
            abort(403)


@app.route('/courses/<int:course_id>/exam/<int:exam_id>')
@login_required
def exam_page(course_id, exam_id):
    user = User.query.filter_by(username=session['user']).first()
    exam, course = _resolve_exam_context(course_id, exam_id)
    _assert_exam_permissions(user, course, exam)

    active_attempt = _active_attempt_for_user(user.id, exam.id)
    past_attempts = ExamAttempt.query.filter(
        ExamAttempt.user_id == user.id,
        ExamAttempt.exam_id == exam.id,
        ExamAttempt.status != 'in-progress'
    ).order_by(ExamAttempt.start_time.desc()).all()

    return render_template(
        'exam.html',
        exam=exam,
        course=course,
        user=user,
        exam_payload=_serialize_exam(exam),
        active_attempt_id=active_attempt.id if active_attempt else None,
        active_time_remaining=_time_remaining_seconds(active_attempt) if active_attempt else None,
        past_attempts=[_summarize_exam_attempt(attempt) for attempt in past_attempts]
    )


@app.route('/courses/<int:course_id>/exam/<int:exam_id>/start', methods=['POST'])
@login_required
def start_exam(course_id, exam_id):
    user = User.query.filter_by(username=session['user']).first()
    exam, course = _resolve_exam_context(course_id, exam_id)
    _assert_exam_permissions(user, course, exam)

    if not exam.allow_retakes and _user_has_passed_exam(user.id, exam.id):
        return jsonify({
            'error': 'You have already passed this exam and retakes are not allowed.'
        }), 403

    if not exam.allow_retakes:
        completed_attempt = ExamAttempt.query.filter(
            ExamAttempt.user_id == user.id,
            ExamAttempt.exam_id == exam.id,
            ExamAttempt.status != 'in-progress'
        ).order_by(ExamAttempt.id.desc()).first()
        if completed_attempt and completed_attempt.passed is not None:
            return jsonify({
                'error': 'Retakes are disabled for this exam.'
            }), 403

    attempt = _start_exam_attempt(user, exam)
    attempt.start_time = _ensure_utc(attempt.start_time) or utcnow()
    db.session.commit()

    return jsonify({
        'attempt_id': attempt.id,
        'attempt_number': attempt.attempt_number,
        'started_at': attempt.start_time.isoformat() if attempt.start_time else None,
        'time_remaining_seconds': _time_remaining_seconds(attempt),
        'exam': _serialize_exam(exam),
        'autosave_payload': attempt.autosave_payload or {},
        'status': attempt.status,
        'passed': attempt.passed
    })


@app.route('/courses/<int:course_id>/exam/<int:exam_id>/autosave', methods=['POST'])
@login_required
def autosave_exam(course_id, exam_id):
    user = User.query.filter_by(username=session['user']).first()
    exam, course = _resolve_exam_context(course_id, exam_id)
    _assert_exam_permissions(user, course, exam)

    payload = request.get_json(silent=True) or {}
    attempt_id = payload.get('attempt_id')
    responses = payload.get('responses')

    if not attempt_id:
        return jsonify({'error': 'Attempt ID is required.'}), 400

    attempt = ExamAttempt.query.get_or_404(attempt_id)
    if attempt.user_id != user.id or attempt.exam_id != exam.id:
        abort(403)

    if attempt.status != 'in-progress':
        return jsonify({'error': 'This attempt is no longer active.'}), 400

    if _time_remaining_seconds(attempt) <= 0:
        _finish_attempt(attempt)
        db.session.commit()
        return jsonify({'error': 'Time expired. Attempt has been auto-submitted.'}), 410

    attempt.autosave_payload = responses or {}
    db.session.commit()

    return jsonify({'success': True, 'time_remaining_seconds': _time_remaining_seconds(attempt)})


@app.route('/courses/<int:course_id>/exam/<int:exam_id>/submit', methods=['POST'])
@login_required
def submit_exam(course_id, exam_id):
    user = User.query.filter_by(username=session['user']).first()
    exam, course = _resolve_exam_context(course_id, exam_id)
    _assert_exam_permissions(user, course, exam)

    payload = request.get_json(silent=True) or {}
    attempt_id = payload.get('attempt_id')
    responses = payload.get('responses') or []

    if not attempt_id:
        return jsonify({'error': 'Attempt ID is required.'}), 400

    attempt = ExamAttempt.query.get_or_404(attempt_id)

    if attempt.user_id != user.id or attempt.exam_id != exam.id:
        abort(403)

    if attempt.status != 'in-progress':
        return jsonify({'error': 'This attempt has already been submitted.'}), 400

    attempt.start_time = _ensure_utc(attempt.start_time) or utcnow()
    now = utcnow()
    deadline = attempt.start_time + timedelta(minutes=exam.duration_minutes) if attempt.start_time else None
    expired = False
    if deadline:
        expired = now > (deadline + timedelta(seconds=2))

    if expired:
        responses = attempt.autosave_payload or []

    responses_iterable = responses
    if isinstance(responses_iterable, dict):
        responses_iterable = [
            {
                'question_id': int(question_id),
                'response': value
            }
            for question_id, value in responses_iterable.items()
        ]

    question_lookup = {question.id: question for question in exam.questions}

    # Clear previous draft answers
    for existing in list(attempt.answers):
        db.session.delete(existing)
    db.session.flush()

    for item in responses_iterable or []:
        question_id = item.get('question_id') if isinstance(item, dict) else None
        response_data = item.get('response') if isinstance(item, dict) else None

        if response_data is None and isinstance(item, dict):
            if 'selected' in item:
                response_data = {'selected': item.get('selected')}
            elif 'values' in item:
                response_data = {'selected': item.get('values')}
            elif 'text' in item:
                response_data = {'text': item.get('text')}
            elif 'value' in item:
                response_data = item.get('value')

        if not question_id or question_id not in question_lookup:
            continue

        question = question_lookup[question_id]
        grading_mode = _exam_grading_mode(exam)
        auto_pass_subjective = (grading_mode == 'automatic')
        is_correct, points_awarded, normalized = _grade_objective_response(question, response_data, auto_pass_subjective)
        if grading_mode == 'manual':
            is_correct = None
            points_awarded = 0.0

        answer = ExamAnswer(
            attempt_id=attempt.id,
            question_id=question_id,
            response_data=normalized,
            is_correct=is_correct,
            points_awarded=points_awarded if is_correct is not None else 0.0
        )
        db.session.add(answer)

    attempt.autosave_payload = None
    _finish_attempt(attempt, submitted_at=now)
    db.session.commit()

    summary = _summarize_exam_attempt(attempt)
    summary['time_remaining_seconds'] = _time_remaining_seconds(attempt)

    return jsonify({'success': True, 'attempt': summary})


@app.route('/courses/<int:course_id>/exam/<int:exam_id>/status')
@login_required
def exam_status(course_id, exam_id):
    user = User.query.filter_by(username=session['user']).first()
    exam, course = _resolve_exam_context(course_id, exam_id)
    _assert_exam_permissions(user, course, exam)

    attempt_id = request.args.get('attempt_id', type=int)

    if attempt_id:
        attempt = ExamAttempt.query.get_or_404(attempt_id)
        if attempt.user_id != user.id or attempt.exam_id != exam.id:
            abort(403)
    else:
        attempt = _active_attempt_for_user(user.id, exam.id)

    if not attempt:
        return jsonify({'attempt': None})

    if attempt.status == 'in-progress' and _time_remaining_seconds(attempt) <= 0:
        _finish_attempt(attempt)
        db.session.commit()

    summary = _summarize_exam_attempt(attempt)
    summary['time_remaining_seconds'] = _time_remaining_seconds(attempt)

    return jsonify({'attempt': summary})


@app.route('/courses/<int:course_id>/exam/<int:exam_id>/results/<int:attempt_id>')
@login_required
def exam_results(course_id, exam_id, attempt_id):
    user = User.query.filter_by(username=session['user']).first()
    exam, course = _resolve_exam_context(course_id, exam_id)
    _assert_exam_permissions(user, course, exam)

    attempt = ExamAttempt.query.get_or_404(attempt_id)
    if attempt.user_id != user.id or attempt.exam_id != exam.id:
        abort(403)

    summary = _summarize_exam_attempt(attempt)

    return render_template(
        'exam_results.html',
        exam=exam,
        course=course,
        attempt=attempt,
        summary=summary,
        user=user
    )


@app.route('/course/<string:course_name>', defaults={'year': 1})
@app.route('/course/<string:course_name>/<int:year>')
@login_required
def course_page(course_name, year):
    user = User.query.filter_by(username=session['user']).first()

    # Find course & lessons
    course = Course.query.filter_by(name=course_name, year=year).first_or_404()

    allowed_courses = set(user.get_courses())

    # ✅ NEW: Check access via new hierarchy system or legacy system
    has_access = course.user_has_access(user.id) or course_name in allowed_courses

    if not has_access and user.role != 'admin':
        flash("⚠️ You don't have access to this course.")
        return redirect(url_for('courses_dashboard'))

    lessons = Lesson.query.filter_by(course_id=course.id).order_by(Lesson.week).all()

    # Fetch published sub-courses for this course
    child_courses = Course.query.filter_by(parent_id=course.id, is_published=True).order_by(Course.order_index, Course.name).all()
    sub_course_payload = []
    for child in child_courses:
        child_access = child.user_has_access(user.id) or child.name in allowed_courses
        child_free = child.is_free()

        child_progress_record = None
        child_progress_percent = 0
        child_started = False
        if child_access or child_free:
            child_progress_record = UserCourseProgress.query.filter_by(
                user_id=user.id,
                course_id=child.id
            ).first()
            if child_progress_record:
                child_started = child_progress_record.progress > 0
                total_child_lessons = Lesson.query.filter_by(course_id=child.id).count()
                if total_child_lessons:
                    child_progress_percent = int((child_progress_record.progress / total_child_lessons) * 100)

        sub_course_payload.append({
            'course': child,
            'has_access': child_access,
            'is_free': child_free,
            'requires_purchase': bool(child.price and child.price > 0),
            'is_started': child_started,
            'progress_percent': child_progress_percent
        })

    exams = sorted(course.exams, key=lambda ex: (
        ex.trigger_lesson.week if ex.trigger_lesson else 999,
        ex.id
    ))

    progress_record = UserCourseProgress.query.filter_by(
        user_id=user.id,
        course_id=course.id
    ).first()
    progress = progress_record.progress if progress_record else 1

    open_lesson_id = request.args.get('open', type=int)

    quiz_feedback = session.get('quiz_feedback')
    retake_id = request.args.get('retake', type=int)
    if retake_id and quiz_feedback and quiz_feedback.get('lesson_id') == retake_id:
        session.pop('quiz_feedback', None)
        quiz_feedback = None
    else:
        quiz_feedback = session.pop('quiz_feedback', None)

    history_raw = session.get('quiz_history', {})
    history_map = {int(k): v for k, v in history_raw.items() if v.get('passed')}

    agreement = CourseAgreement.query.filter_by(user_id=user.id, course_id=course.id).first()

    exam_states = []
    exams_by_trigger = {}
    locked_weeks = set()
    locking_exam_by_week = {}

    for exam in exams:
        latest_attempt = ExamAttempt.query.filter_by(user_id=user.id, exam_id=exam.id).order_by(ExamAttempt.start_time.desc()).first()
        in_progress = latest_attempt and latest_attempt.status == 'in-progress'
        settings = exam.settings or {}
        unlock_on_submission = bool(settings.get('unlock_on_submission'))
        passed = bool(latest_attempt and latest_attempt.passed)
        latest_status = latest_attempt.status if latest_attempt else 'not_started'
        submitted_unlock = unlock_on_submission and latest_attempt and latest_status in {'submitted', 'graded', 'passed'}
        progress_unlocked = passed or submitted_unlock
        exam_state = {
            'exam': exam,
            'latest_attempt': latest_attempt,
            'passed': passed,
            'status': latest_status,
            'resume_attempt_id': latest_attempt.id if in_progress else None,
            'time_remaining': _time_remaining_seconds(latest_attempt) if in_progress else None,
            'trigger_lesson_id': exam.trigger_lesson_id,
            'unlock_on_submission': unlock_on_submission,
            'progress_unlocked': progress_unlocked
        }
        exam_states.append(exam_state)

        trigger_key = exam.trigger_lesson_id or 0
        exams_by_trigger.setdefault(trigger_key, []).append(exam_state)

        if exam.is_required and not progress_unlocked:
            trigger_week = exam.trigger_lesson.week if exam.trigger_lesson else (lessons[-1].week if lessons else 0)
            for lesson in lessons:
                if lesson.week > trigger_week:
                    locked_weeks.add(lesson.week)
                    locking_exam_by_week.setdefault(lesson.week, exam_state)

    return render_template(
        "courses.html",
        current_user=user,
        course=course,
        lessons=lessons,
        progress=progress,
        open_lesson_id=open_lesson_id,
        quiz_feedback=quiz_feedback,
        quiz_history=history_map,
        course_agreement=agreement,
        exams=exam_states,
        exams_by_trigger=exams_by_trigger,
        locked_weeks=locked_weeks,
        locking_exam_by_week=locking_exam_by_week,
        sub_courses=sub_course_payload
    )


@app.route('/course/<int:course_id>')
@login_required
def course_detail(course_id):
    """
    View course by ID - supports new hierarchy system
    """
    user = User.query.filter_by(username=session['user']).first()
    course = Course.query.get_or_404(course_id)

    # ✅ Check access via new hierarchy system or legacy system
    has_access = course.user_has_access(user.id) or course.name in user.get_courses()

    if not has_access and user.role != 'admin':
        flash("⚠️ You don't have access to this course.")
        return redirect(url_for('courses_dashboard'))

    # Redirect to the existing course_page route with name and year
    return redirect(url_for('course_page', course_name=course.name, year=course.year))


# -------------------- Stripe Payment Routes -------------------- #
@app.route('/courses/<int:course_id>/checkout')
@login_required
def course_checkout(course_id):
    """
    Create a Stripe Checkout session for purchasing a course.
    """
    from stripe_helpers import create_checkout_session, sync_course_with_stripe

    user = User.query.filter_by(username=session['user']).first()
    course = Course.query.get_or_404(course_id)

    # Check if course is free
    if not course.price or course.price <= 0:
        flash("⚠️ This course is free. No payment required.")
        return redirect(url_for('courses_dashboard'))

    # Check if user already has access
    if course.user_has_access(user.id):
        flash("✅ You already have access to this course.")
        return redirect(url_for('course_page', course_name=course.name, year=course.year))

    # Sync course with Stripe if not already synced
    if not course.stripe_price_id:
        try:
            product_id, price_id = sync_course_with_stripe(course)
            course.stripe_product_id = product_id
            course.stripe_price_id = price_id
            db.session.commit()
        except Exception as e:
            app.logger.error(f"Failed to sync course {course_id} with Stripe: {e}")
            flash("⚠️ Payment system error. Please contact support.")
            return redirect(url_for('courses_dashboard'))

    # Create checkout session
    try:
        success_url = url_for('payment_success', _external=True) + f'?session_id={{CHECKOUT_SESSION_ID}}'
        cancel_url = url_for('payment_cancel', course_id=course_id, _external=True)

        checkout_session = create_checkout_session(
            course=course,
            user=user,
            success_url=success_url,
            cancel_url=cancel_url
        )

        return redirect(checkout_session.url, code=303)

    except Exception as e:
        app.logger.error(f"Failed to create checkout session: {e}")
        flash("⚠️ Could not create payment session. Please try again.")
        return redirect(url_for('courses_dashboard'))


@app.route('/payment/success')
@login_required
def payment_success():
    """
    Handle successful payment - shown after Stripe checkout.
    Note: Access is granted via webhook for security, not here.
    """
    session_id = request.args.get('session_id')
    return render_template('payment_success.html', session_id=session_id)


@app.route('/payment/cancel/<int:course_id>')
@login_required
def payment_cancel(course_id):
    """
    Handle cancelled payment.
    """
    course = Course.query.get_or_404(course_id)
    return render_template('payment_cancel.html', course=course)


# -------------------- Subscription Routes -------------------- #
@app.route('/subscriptions')
@login_required
def subscription_plans():
    """
    Show available subscription plans.
    """
    plans = SubscriptionPlan.query.filter_by(is_active=True).all()
    user = User.query.filter_by(username=session['user']).first()
    active_subscription = get_user_active_subscription(user.id) if user else None

    return render_template('subscription_plans.html', plans=plans, active_subscription=active_subscription)


@app.route('/subscribe/<int:plan_id>')
@login_required
def subscribe_plan(plan_id):
    """
    Create a Stripe Checkout session for subscription.
    """
    from stripe_helpers import create_subscription_checkout_session, sync_plan_with_stripe

    user = User.query.filter_by(username=session['user']).first()
    plan = SubscriptionPlan.query.get_or_404(plan_id)

    if not plan.is_active:
        flash("⚠️ This subscription plan is not available.")
        return redirect(url_for('subscription_plans'))

    # Check if user already has an active subscription
    active_sub = get_user_active_subscription(user.id)
    if active_sub:
        flash("⚠️ You already have an active subscription.")
        return redirect(url_for('my_subscription'))

    # Sync plan with Stripe if not already synced
    if not plan.stripe_price_id:
        try:
            product_id, price_id = sync_plan_with_stripe(plan)
            plan.stripe_product_id = product_id
            plan.stripe_price_id = price_id
            db.session.commit()
        except Exception as e:
            app.logger.error(f"Failed to sync plan {plan_id} with Stripe: {e}")
            flash("⚠️ Payment system error. Please contact support.")
            return redirect(url_for('subscription_plans'))

    # Create checkout session
    try:
        success_url = url_for('subscription_success', _external=True) + '?session_id={CHECKOUT_SESSION_ID}'
        cancel_url = url_for('subscription_cancel', plan_id=plan_id, _external=True)

        checkout_session = create_subscription_checkout_session(
            plan=plan,
            user=user,
            success_url=success_url,
            cancel_url=cancel_url
        )

        return redirect(checkout_session.url, code=303)

    except Exception as e:
        app.logger.error(f"Failed to create subscription checkout: {e}")
        flash("⚠️ Could not create subscription. Please try again.")
        return redirect(url_for('subscription_plans'))


@app.route('/subscription/success')
@login_required
def subscription_success():
    """
    Handle successful subscription - shown after Stripe checkout.
    """
    session_id = request.args.get('session_id')
    return render_template('subscription_success.html', session_id=session_id)


@app.route('/subscription/cancel/<int:plan_id>')
@login_required
def subscription_cancel(plan_id):
    """
    Handle cancelled subscription checkout.
    """
    plan = SubscriptionPlan.query.get_or_404(plan_id)
    return render_template('subscription_cancel.html', plan=plan)


@app.route('/my-subscription')
@login_required
def my_subscription():
    """
    Show user's subscription details and management options.
    """
    user = User.query.filter_by(username=session['user']).first()
    active_subscription = get_user_active_subscription(user.id)

    subscription_data = None
    if active_subscription:
        plan = SubscriptionPlan.query.get(active_subscription.plan_id)
        courses = plan.get_courses() if plan else []

        subscription_data = {
            'subscription': active_subscription,
            'plan': plan,
            'courses': courses
        }

    # Get all subscriptions (including canceled)
    all_subscriptions = UserSubscription.query.filter_by(user_id=user.id).order_by(UserSubscription.start_date.desc()).all()

    return render_template('my_subscription.html', subscription_data=subscription_data, all_subscriptions=all_subscriptions)


@app.route('/billing-portal')
@login_required
def billing_portal():
    """
    Redirect to Stripe Customer Portal for subscription management.
    """
    from stripe_helpers import create_customer_portal_session

    user = User.query.filter_by(username=session['user']).first()
    active_subscription = get_user_active_subscription(user.id)

    if not active_subscription or not active_subscription.stripe_customer_id:
        flash("⚠️ No active subscription found.")
        return redirect(url_for('my_subscription'))

    try:
        return_url = url_for('my_subscription', _external=True)
        portal_session = create_customer_portal_session(
            active_subscription.stripe_customer_id,
            return_url
        )
        return redirect(portal_session.url, code=303)

    except Exception as e:
        app.logger.error(f"Failed to create billing portal session: {e}")
        flash("⚠️ Could not access billing portal. Please try again.")
        return redirect(url_for('my_subscription'))


@app.route('/stripe/webhook', methods=['POST'])
def stripe_webhook():
    """
    Handle Stripe webhook events.
    This is called by Stripe to notify us of payment events.
    """
    from stripe_helpers import (
        verify_webhook_signature, handle_checkout_session_completed,
        handle_subscription_created, handle_invoice_payment_succeeded,
        handle_invoice_payment_failed
    )
    import stripe

    payload = request.get_data()
    sig_header = request.headers.get('Stripe-Signature')

    try:
        event = verify_webhook_signature(payload, sig_header)
    except ValueError:
        # Invalid payload
        app.logger.error("Invalid webhook payload")
        return jsonify({'error': 'Invalid payload'}), 400
    except stripe.error.SignatureVerificationError:
        # Invalid signature
        app.logger.error("Invalid webhook signature")
        return jsonify({'error': 'Invalid signature'}), 400
    except Exception as e:
        app.logger.error(f"Webhook verification error: {e}")
        return jsonify({'error': 'Verification failed'}), 400

    # Handle the event
    event_type = event['type']

    if event_type == 'checkout.session.completed':
        session = event['data']['object']
        mode = session.get('mode')

        if mode == 'subscription':
            # Handle subscription checkout
            metadata = session.get('metadata', {})
            subscription_id = session.get('subscription')

            if subscription_id:
                try:
                    # Fetch full subscription details from Stripe
                    subscription = stripe.Subscription.retrieve(subscription_id)
                    sub_data = handle_subscription_created(subscription)

                    # Create or update user subscription record
                    user_sub = UserSubscription.query.filter_by(
                        stripe_subscription_id=subscription_id
                    ).first()

                    if not user_sub:
                        user_sub = UserSubscription(
                            user_id=sub_data['user_id'],
                            plan_id=sub_data['plan_id'],
                            stripe_subscription_id=sub_data['subscription_id'],
                            stripe_customer_id=sub_data['customer_id'],
                            status=sub_data['status'],
                            start_date=datetime.fromtimestamp(sub_data['current_period_start']),
                            current_period_start=datetime.fromtimestamp(sub_data['current_period_start']),
                            current_period_end=datetime.fromtimestamp(sub_data['current_period_end'])
                        )

                        if sub_data['trial_start']:
                            user_sub.trial_start = datetime.fromtimestamp(sub_data['trial_start'])
                        if sub_data['trial_end']:
                            user_sub.trial_end = datetime.fromtimestamp(sub_data['trial_end'])

                        db.session.add(user_sub)

                    # Grant access to courses in the plan
                    plan = SubscriptionPlan.query.get(sub_data['plan_id'])
                    if plan and plan.course_ids:
                        grant_subscription_access(sub_data['user_id'], plan.course_ids)

                    # Upgrade user role
                    user = User.query.get(sub_data['user_id'])
                    if user and user.role == 'user':
                        user.role = 'subscriber'

                    db.session.commit()
                    app.logger.info(f"Created subscription {subscription_id} for user {sub_data['user_id']}")

                except Exception as e:
                    app.logger.error(f"Failed to create subscription: {e}")
                    db.session.rollback()
                    return jsonify({'error': 'Failed to create subscription'}), 500
        else:
            # Handle one-time purchase
            payment_data = handle_checkout_session_completed(session)

            try:
                user = User.query.get(payment_data['user_id'])
                course = Course.query.get(payment_data['course_id'])

                if user and course:
                    grant_course_access(
                        user.id,
                        course,
                        access_type='purchased',
                        amount_paid=payment_data['amount_paid'],
                        payment_intent_id=payment_data['payment_intent_id']
                    )

                    if user.role == 'user':
                        user.role = 'paid'

                    db.session.commit()
                    app.logger.info(f"Granted access to course {course.id} for user {user.id}")

                    # Send course enrollment email
                    try:
                        from email_utils import send_course_enrollment_email
                        send_course_enrollment_email(user, course.name, payment_data['amount_paid'])
                    except Exception as e:
                        app.logger.error(f"Failed to send enrollment email: {e}")

            except Exception as e:
                app.logger.error(f"Failed to grant course access: {e}")
                db.session.rollback()
                return jsonify({'error': 'Failed to grant access'}), 500

    elif event_type == 'invoice.payment_succeeded':
        invoice = event['data']['object']
        invoice_data = handle_invoice_payment_succeeded(invoice)

        try:
            user_sub = UserSubscription.query.filter_by(
                stripe_subscription_id=invoice_data['subscription_id']
            ).first()

            if user_sub:
                user_sub.status = 'active'
                user_sub.last_payment_date = datetime.utcnow()
                user_sub.last_payment_amount = invoice_data['amount_paid']
                user_sub.current_period_start = datetime.fromtimestamp(invoice_data['period_start'])
                user_sub.current_period_end = datetime.fromtimestamp(invoice_data['period_end'])

                # Unlock courses if they were locked
                plan = SubscriptionPlan.query.get(user_sub.plan_id)
                if plan and plan.course_ids:
                    unlock_subscription_courses(user_sub.user_id, plan.course_ids)

                db.session.commit()
                app.logger.info(f"Payment succeeded for subscription {invoice_data['subscription_id']}")

        except Exception as e:
            app.logger.error(f"Failed to update subscription on payment success: {e}")
            db.session.rollback()

    elif event_type == 'invoice.payment_failed':
        invoice = event['data']['object']
        invoice_data = handle_invoice_payment_failed(invoice)

        try:
            user_sub = UserSubscription.query.filter_by(
                stripe_subscription_id=invoice_data['subscription_id']
            ).first()

            if user_sub:
                user_sub.status = 'past_due'
                plan = SubscriptionPlan.query.get(user_sub.plan_id)

                # Set grace period end date
                if plan:
                    user_sub.end_date = datetime.utcnow() + timedelta(days=plan.grace_period_days)

                db.session.commit()
                app.logger.info(f"Payment failed for subscription {invoice_data['subscription_id']}, entering grace period")

                # TODO: Send email notification

        except Exception as e:
            app.logger.error(f"Failed to update subscription on payment failure: {e}")
            db.session.rollback()

    elif event_type == 'customer.subscription.updated':
        subscription = event['data']['object']

        try:
            user_sub = UserSubscription.query.filter_by(
                stripe_subscription_id=subscription['id']
            ).first()

            if user_sub:
                user_sub.status = subscription['status']
                user_sub.cancel_at_period_end = subscription.get('cancel_at_period_end', False)

                if subscription.get('canceled_at'):
                    user_sub.canceled_at = datetime.fromtimestamp(subscription['canceled_at'])

                db.session.commit()
                app.logger.info(f"Updated subscription {subscription['id']}")

        except Exception as e:
            app.logger.error(f"Failed to update subscription: {e}")
            db.session.rollback()

    elif event_type == 'customer.subscription.deleted':
        subscription = event['data']['object']

        try:
            user_sub = UserSubscription.query.filter_by(
                stripe_subscription_id=subscription['id']
            ).first()

            if user_sub:
                user_sub.status = 'canceled'
                user_sub.end_date = datetime.utcnow()

                # Lock course access but retain data
                plan = SubscriptionPlan.query.get(user_sub.plan_id)
                if plan and plan.course_ids:
                    revoke_subscription_access(user_sub.user_id, plan.course_ids, keep_progress=True)

                db.session.commit()
                app.logger.info(f"Canceled subscription {subscription['id']}, data retained")

        except Exception as e:
            app.logger.error(f"Failed to cancel subscription: {e}")
            db.session.rollback()

    return jsonify({'status': 'success'}), 200


# -------------------- Error Handlers -------------------- #
@app.errorhandler(500)
def internal_error(error):
    """Handle 500 Internal Server Error"""
    import traceback
    error_details = traceback.format_exc()
    app.logger.error(f'Internal Server Error: {error_details}')
    if app.debug:
        return f"<h1>Internal Server Error</h1><pre>{error_details}</pre>", 500
    return render_template('error.html', error_code=500, error_message="Internal Server Error"), 500


@app.errorhandler(404)
def not_found_error(error):
    """Handle 404 Not Found Error"""
    return render_template('error.html', error_code=404, error_message="Page Not Found"), 404


@app.errorhandler(403)
def forbidden_error(error):
    """Handle 403 Forbidden Error"""
    return render_template('error.html', error_code=403, error_message="Access Forbidden"), 403


if __name__ == '__main__':
    # Get host and port from environment variables
    host = os.environ.get('HOST', '0.0.0.0')
    port = int(os.environ.get('PORT', 5005))
    debug = os.environ.get('DEBUG', 'False').lower() == 'true'

    app.run(host=host, port=port, debug=debug)
